#!/usr/bin/env python3
"""
视频合成工具 Pro 版
支持：字幕样式自定义、多种转场、片头片尾、批量处理

使用方法:
    python3 video_generator_pro.py --project 项目路径 [选项]

示例:
    # 基础用法（自动使用 04_videos 或 05_audio + 图片）
    python3 video_generator_pro.py --project projects/XXX

    # 带字幕 + 片头片尾 + 转场
    python3 video_generator_pro.py --project projects/XXX \
        --subtitle --subtitle-style news \
        --transition fade \
        --intro intro.mp4 --outro outro.mp4

    # 批量处理
    python3 video_generator_pro.py --batch projects/*
"""

import os
import sys
import json
import subprocess
import argparse
import asyncio
import re
import time

try:
    import readline
except ImportError:
    pass
import datetime
import concurrent.futures
import logging
import traceback
from pathlib import Path
from typing import List, Tuple, Optional, Dict
import tempfile
import shutil
from dataclasses import dataclass

try:
    import edge_tts
    EDGE_TTS_AVAILABLE = True
except ImportError:
    EDGE_TTS_AVAILABLE = False


# 字幕样式预设
SUBTITLE_STYLES = {
    'news': {
        'font': 'SourceHanSansSC-Bold',
        'size': 48,
        'color': 'white',
        'box': 1,
        'boxcolor': 'black@0.7',
        'boxborderw': 10,
        'borderw': 2,
        'bordercolor': 'black',
        'position': 'bottom',
        'y_offset': 100
    },
    'youtube': {
        'font': 'Arial-Bold',
        'size': 42,
        'color': 'yellow',
        'box': 0,
        'borderw': 3,
        'bordercolor': 'black',
        'position': 'bottom',
        'y_offset': 80
    },
    'minimal': {
        'font': 'SourceHanSansSC-Regular',
        'size': 40,
        'color': 'white',
        'box': 0,
        'borderw': 2,
        'bordercolor': 'black@0.5',
        'position': 'center',
        'y_offset': 0
    },
    'tiktok': {
        'font': 'SourceHanSansSC-Bold',
        'size': 52,
        'color': 'white',
        'box': 1,
        'boxcolor': 'black@0.5',
        'boxborderw': 8,
        'borderw': 3,
        'bordercolor': 'black',
        'position': 'center',
        'y_offset': 100
    }
}


# 转场效果预设（基于 ffmpeg xfade 滤镜，需 ffmpeg 4.3+）
TRANSITIONS = {
    'fade': 'fade',
    'dissolve': 'dissolve',
    'wipeleft': 'wipeleft',
    'wiperight': 'wiperight',
    'wipeup': 'wipeup',
    'wipedown': 'wipedown',
    'slideleft': 'slideleft',
    'slideright': 'slideright',
    'slideup': 'slideup',
    'slidedown': 'slidedown',
    'smoothleft': 'smoothleft',
    'smoothright': 'smoothright',
    'smoothup': 'smoothup',
    'smoothdown': 'smoothdown',
    'circlecrop': 'circlecrop',
    'rectcrop': 'rectcrop',
    'circleclose': 'circleclose',
    'circleopen': 'circleopen',
    'horzclose': 'horzclose',
    'horzopen': 'horzopen',
    'vertclose': 'vertclose',
    'vertopen': 'vertopen',
    'diagbl': 'diagbl',
    'diagbr': 'diagbr',
    'diagtl': 'diagtl',
    'diagtr': 'diagtr',
    'hlslice': 'hlslice',
    'hrslice': 'hrslice',
    'vuslice': 'vuslice',
    'vdslice': 'vdslice',
    'pixelize': 'pixelize',
    'radial': 'radial',
    'distance': 'distance',
    'fadeblack': 'fadeblack',
    'fadewhite': 'fadewhite',
    'none': None
}


# 项目模板预设
PROJECT_TEMPLATES = {
    'news': {
        'mode': 'image',
        'resolution': '1920x1080',
        'fps': 30,
        'subtitle_style': 'news',
        'transition': 'fade',
        'voice': 'Yunjian',
        'article': """# 今日新闻

@全局:云健
@默认图: 新闻台标

欢迎收看今日新闻。

@女声: @图:现场 记者现场报道，投资者情绪高涨。

@云健: 第二条新闻：国际油价回落。

@女声: @图:油价 布伦特原油跌至80美元。

@云健: 第三条新闻：央行宣布降准。

感谢收看今日新闻。
"""
    },
    'food': {
        'mode': 'image',
        'resolution': '1080x1920',
        'fps': 30,
        'subtitle_style': 'tiktok',
        'transition': 'pixelize',
        'voice': 'Xiaoxiao',
        'article': """# 美食探店

@全局:女声
@默认图: 餐厅环境

@男声: 大家好，欢迎来到美食探店。

@图:招牌菜
@男声: 这是我们的招牌菜，色香味俱全。

回到女声介绍餐厅环境。

@图:厨师
@男声: 我们的大厨有20年经验。

@图:食材
@女声: 所有食材都是当天采购。

@图:价格
@男声: 价格非常亲民。

@女声: 欢迎来品尝！
"""
    },
    'tutorial': {
        'mode': 'image',
        'resolution': '1920x1080',
        'fps': 30,
        'subtitle_style': 'youtube',
        'transition': 'slideleft',
        'voice': 'Xiaoyi',
        'article': """# Python入门课程

@全局:晓伊
@默认图: 课程封面

欢迎来到Python入门课程。我是主讲老师李老师。

@云希: 我是助教小张，负责解答大家的问题。

今天我们先了解什么是Python。

@云希: Python是一种简单易学但功能强大的编程语言。

接下来我们安装Python环境。
"""
    },
    'education': {
        'mode': 'image',
        'resolution': '1920x1080',
        'fps': 30,
        'subtitle_style': 'minimal',
        'transition': 'fade',
        'voice': 'Yunyang',
        'article': """# 历史知识科普

@全局:男声
@默认图: 历史背景

今天我们来聊聊三国时期的故事。

@图:诸葛亮
@女声: 诸葛亮是蜀汉的丞相，以智慧著称。

@男声: 他草船借箭、空城计的故事家喻户晓。

@图:赤壁
@女声: 赤壁之战是三国时期最著名的战役之一。

@男声: 感谢收听，我们下期再见！
"""
    }
}


@dataclass
class Scene:
    """场景数据类"""
    index: int
    audio_path: Optional[Path]
    video_path: Optional[Path]
    image_path: Optional[Path]
    subtitle: str
    duration: float


class Tee:
    """同时输出到终端和文件的 stdout 重定向器"""
    def __init__(self, file_path: Path, encoding='utf-8'):
        self.terminal = sys.stdout
        self.log_file = open(file_path, 'w', encoding=encoding, buffering=1)

    def write(self, message):
        self.terminal.write(message)
        self.log_file.write(message)
        self.log_file.flush()

    def flush(self):
        self.terminal.flush()
        self.log_file.flush()

    def close(self):
        if self.log_file and not self.log_file.closed:
            self.log_file.close()


def setup_logging(project_dir: Path, preview_mode: bool = False) -> Optional[Tee]:
    """设置日志：同时输出到终端和日志文件"""
    final_dir = project_dir / '07_final'
    final_dir.mkdir(exist_ok=True)

    now = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    log_name = f"generate_{now}.log"
    log_path = final_dir / log_name

    tee = Tee(log_path)
    sys.stdout = tee

    print(f"\n📝 运行日志已保存: {log_path}")
    return tee


# ffprobe 时长缓存（避免重复调用）
_DURATION_CACHE: Dict[str, float] = {}


def get_media_duration(path: str) -> float:
    """获取媒体时长（秒），带缓存"""
    # 使用路径+修改时间作为缓存键，文件变更自动失效
    try:
        mtime = str(os.path.getmtime(path))
    except OSError:
        mtime = '0'
    cache_key = f"{path}:{mtime}"

    if cache_key in _DURATION_CACHE:
        return _DURATION_CACHE[cache_key]

    cmd = [
        'ffprobe', '-v', 'error', '-show_entries', 'format=duration',
        '-of', 'default=noprint_wrappers=1:nokey=1', path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    try:
        duration = float(result.stdout.strip())
        _DURATION_CACHE[cache_key] = duration
        return duration
    except:
        _DURATION_CACHE[cache_key] = 0.0
        return 0.0


def run_ffmpeg(cmd: list, max_retries: int = 2, check_output: bool = True) -> subprocess.CompletedProcess:
    """运行 ffmpeg 命令，支持失败重试"""
    for attempt in range(1, max_retries + 1):
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode == 0:
            return result
        if attempt < max_retries:
            stderr_short = result.stderr[:200].replace('\n', ' ')
            print(f"   ⚠️  ffmpeg 失败 (尝试 {attempt}/{max_retries}): {stderr_short}")
            print(f"   ⏳  2秒后重试...")
            time.sleep(2)
        else:
            print(f"   ❌ ffmpeg 最终失败 (已重试 {max_retries} 次)")
            if check_output and result.stderr:
                print(f"   错误详情: {result.stderr[:500]}")
    return result


def wrap_subtitle_text(text: str, max_chars: int = 14) -> str:
    """智能字幕换行：按标点拆分，控制每行字数，仅在行满时换行

    Args:
        text: 原始字幕文本
        max_chars: 每行最大中文字符数（默认14个）

    Returns:
        带\\n换行符的字幕文本
    """
    if not text:
        return text

    import re
    # 按主要标点拆分，保留标点
    segments = re.split(r'([，。！？；：])', text.strip())

    lines = []
    current_line = ''

    for seg in segments:
        if not seg:
            continue
        # 如果是标点，追加到当前行
        if seg in '，。！？；：':
            current_line += seg
            # 标点后不强制换行，继续累积
            continue

        # 纯文本段，检查追加后是否超过 max_chars
        if len(current_line) + len(seg) <= max_chars:
            current_line += seg
        else:
            # 当前行满了，先保存
            if current_line:
                lines.append(current_line)
                current_line = ''
            # 新段如果还是超长，强制按 max_chars 截断
            while len(seg) > max_chars:
                lines.append(seg[:max_chars])
                seg = seg[max_chars:]
            current_line = seg

    if current_line:
        lines.append(current_line)

    return '\\n'.join(lines)


def build_subtitle_filter(
    subtitle: str,
    subtitle_style: Dict,
    width: int,
    height: int,
    animation: str = 'none',
    anim_duration: float = 0.5,
    wrap: bool = True,
    max_chars_per_line: int = 14
) -> str:
    """构建字幕滤镜字符串，支持动画效果和自动换行"""
    # 自动换行处理
    if wrap:
        subtitle = wrap_subtitle_text(subtitle, max_chars=max_chars_per_line)

    # 处理换行符：wrap_subtitle_text 返回的 \n 字面量在某些 ffmpeg 版本中不能正确换行，
    # 先替换为空格；同时把真正的换行符也替换为空格，避免 filter 字符串断裂
    subtitle = subtitle.replace('\\n', ' ').replace('\n', ' ')
    # 转义剩余的反斜杠，保护 filter 语法
    subtitle = subtitle.replace('\\', '\\\\')

    y_pos = height - subtitle_style['y_offset'] if subtitle_style['position'] == 'bottom' else height // 2

    # 基础 drawtext 参数
    base = (
        f"drawtext=fontfile=/System/Library/Fonts/STHeiti\\ Medium.ttc:"
        f"text='{subtitle}':"
        f"fontcolor={subtitle_style['color']}:"
        f"fontsize={subtitle_style['size']}:"
        f"borderw={subtitle_style['borderw']}:"
        f"bordercolor={subtitle_style['bordercolor']}"
    )

    # 动画效果
    if animation == 'slide_up':
        start_y = height + 100
        y_expr = f"if(lt(t\\,{anim_duration})\\,{start_y}-({start_y}-{y_pos})/{anim_duration}*t\\,{y_pos})"
        base += f":y={y_expr}:x=(w-text_w)/2"
    elif animation == 'fade_in':
        alpha_expr = f"if(lt(t\\,{anim_duration})\\,t/{anim_duration}\\,1)"
        base += f":y={y_pos}:x=(w-text_w)/2:alpha={alpha_expr}"
    else:
        base += f":x=(w-text_w)/2:y={y_pos}"

    if subtitle_style.get('box'):
        base += f":box=1:boxcolor={subtitle_style['boxcolor']}:boxborderw={subtitle_style['boxborderw']}"

    return base


def build_sentence_subtitle_filter(
    subtitle: str,
    subtitle_style: Dict,
    width: int,
    height: int,
    total_duration: float,
    animation: str = 'none',
    anim_duration: float = 0.5,
    wrap: bool = True,
    max_chars_per_line: int = 14,
    subtitle_gap: float = 0.1
) -> str:
    """构建逐句字幕滤镜：按句子拆分，根据时长计算每句显示时间

    Args:
        subtitle_gap: 句间黑屏间隔（秒）。0=无缝衔接下一句，>0=读完后黑屏N秒再出下一句（默认0.1秒）
    """
    import re

    # 按句子拆分（支持中文标点）
    sentences = re.split(r'([。！？；])', subtitle.strip())
    # 重组：把标点和前面的内容合并
    items = []
    buf = ''
    for s in sentences:
        if s in '。！？；':
            buf += s
            if buf.strip():
                items.append(buf.strip())
            buf = ''
        else:
            buf += s
    if buf.strip():
        items.append(buf.strip())

    if not items:
        items = [subtitle.strip()]

    # 单句直接走整段显示
    if len(items) == 1:
        return build_subtitle_filter(
            subtitle, subtitle_style, width, height,
            animation=animation, anim_duration=anim_duration,
            wrap=wrap, max_chars_per_line=max_chars_per_line
        )

    # 按字数分配时长
    total_chars = sum(len(s) for s in items)
    char_time = total_duration / total_chars if total_chars > 0 else total_duration

    # 计算每句起止时间
    filters = []
    current_time = 0.0
    y_pos = height - subtitle_style['y_offset'] if subtitle_style['position'] == 'bottom' else height // 2

    for idx, sentence in enumerate(items):
        # 处理换行符和反斜杠，避免 ffmpeg filter 解析出错
        sentence = sentence.replace('\n', ' ').replace('\\n', ' ')
        sentence = sentence.replace("\\", "\\\\").replace("'", "'\\''")
        sent_chars = len(sentence)
        sent_duration = max(sent_chars * char_time, 1.0)  # 最少1秒

        start_t = current_time
        # 非最后一句：读完之后可配置间隔（subtitle_gap）再出下一句
        if idx == len(items) - 1:
            end_t = total_duration
        else:
            end_t = min(current_time + sent_duration - subtitle_gap, total_duration)
            # 确保 end_t 不会早于 start_t + 0.3（每句至少显示0.3秒）
            if end_t < start_t + 0.3:
                end_t = start_t + 0.3

        # 构建单句 drawtext，带 enable 时间窗口
        ft = (
            f"drawtext=fontfile=/System/Library/Fonts/STHeiti\\ Medium.ttc:"
            f"text='{sentence}':"
            f"fontcolor={subtitle_style['color']}:"
            f"fontsize={subtitle_style['size']}:"
            f"borderw={subtitle_style['borderw']}:"
            f"bordercolor={subtitle_style['bordercolor']}:"
            f"x=(w-text_w)/2:y={y_pos}:"
            f"enable='between(t\\,{start_t:.2f}\\,{end_t:.2f})'"
        )

        if subtitle_style.get('box'):
            ft += f":box=1:boxcolor={subtitle_style['boxcolor']}:boxborderw={subtitle_style['boxborderw']}"

        filters.append(ft)
        current_time += sent_duration

    return ','.join(filters)


def generate_text_video(
    text: str,
    output_path: Path,
    resolution: Tuple[int, int] = (1920, 1080),
    duration: float = 3.0,
    fps: int = 30,
    bg_color: str = 'black',
    text_color: str = 'white',
    font_size: int = 72
) -> bool:
    """生成文字+纯色背景的视频（用于片头/片尾模板），带静音音轨确保拼接兼容"""
    width, height = resolution

    # 安全处理文本中的特殊字符（单引号等）
    safe_text = text.replace("'", "'\\''")

    cmd = [
        'ffmpeg', '-y',
        '-f', 'lavfi', '-i', f"color=c={bg_color}:s={width}x{height}:d={duration}:r={fps}",
        '-f', 'lavfi', '-i', f"anullsrc=r=48000:cl=stereo",  # 静音音轨
        '-vf', (
            f"drawtext=fontfile=/System/Library/Fonts/STHeiti\\ Medium.ttc:"
            f"text='{safe_text}':fontcolor={text_color}:fontsize={font_size}:"
            f"x=(w-text_w)/2:y=(h-text_h)/2:borderw=4:bordercolor=black:"
            f"alpha='if(lt(t,0.3),t/0.3,if(lt(t,{duration-0.3}),1,({duration}-t)/0.3))'"
        ),
        '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
        '-pix_fmt', 'yuv420p',
        '-c:a', 'aac', '-b:a', '128k',
        '-r', str(fps),
        '-shortest',
        '-t', str(duration),
        str(output_path)
    ]

    try:
        result = run_ffmpeg(cmd, max_retries=1, check_output=False)
        if result.returncode == 0 and output_path.exists():
            print(f"   ✅ 文字视频生成: {output_path.name} ({duration:.1f}s)")
            return True
        print(f"   ⚠️  文字视频生成失败")
        return False
    except Exception as e:
        print(f"   ⚠️  文字视频生成异常: {e}")
        return False


def transcribe_video_with_whisper(video_path: Path, output_article: Path, model_size: str = 'small') -> bool:
    """使用本地 faster-whisper 识别视频语音，生成文章

    首次使用会自动从 HuggingFace Hub 下载模型（约 244MB，small 模型）。
    国内网络不稳定时可设置镜像：
        export HF_ENDPOINT=https://hf-mirror.com
    """
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        print("   ❌ 未安装 faster-whisper，请运行: pip install faster-whisper")
        return False

    # 国内用户自动使用 hf-mirror 镜像（如果未手动设置）
    if not os.environ.get('HF_ENDPOINT') and not os.environ.get('HF_HUB_OFFLINE'):
        os.environ.setdefault('HF_ENDPOINT', 'https://hf-mirror.com')

    try:
        from faster_whisper.utils import download_model, _MODELS
        # 先检查模型是否已本地缓存，已缓存时跳过网络检查（秒开）
        repo_id = _MODELS[model_size]
        cache_dir = Path.home() / '.cache' / 'huggingface' / 'hub'
        # 构造快照目录路径（不依赖网络）
        expected_path = cache_dir / f"models--{repo_id.replace('/', '--')}" / "snapshots"
        if expected_path.exists() and any(expected_path.iterdir()):
            # 已缓存，直接取最新快照
            snapshots = sorted(expected_path.iterdir())
            if snapshots:
                model_path = str(snapshots[-1])
                print(f"   ✅ 模型已缓存: {model_size}")
        else:
            # 未缓存，需要下载
            print(f"   📦 下载模型... (模型: {model_size})")
            model_path = download_model(model_size)
            print(f"   ✅ 模型下载完成")
    except Exception as e:
        err_msg = str(e)
        if 'ConnectError' in err_msg or 'SSL' in err_msg or 'huggingface' in err_msg.lower():
            print("   ❌ 模型下载失败: 网络连接异常（HuggingFace 访问受限）")
            print("   💡 解决方案（三选一）:")
            print("      1. 设置镜像后重试: export HF_ENDPOINT=https://hf-mirror.com")
            print("      2. 手动下载模型放到 ~/.cache/huggingface/hub/")
            print("      3. 使用离线模式（已下载过模型）: export HF_HUB_OFFLINE=1")
        else:
            print(f"   ❌ 模型下载失败: {e}")
        return False

    print(f"   🎙️  开始语音识别...")
    try:
        model = WhisperModel(model_size, device="cpu", compute_type="int8")
        segments, info = model.transcribe(str(video_path), language="zh", beam_size=5)

        print(f"   检测语言: {info.language} (概率: {info.language_probability:.2f})")

        # 收集所有文本，合并为一段（视频模式通常单个视频配完整字幕）
        texts = []
        for segment in segments:
            texts.append(segment.text.strip())

        full_text = ' '.join(texts)

        # 保存文章
        output_article.parent.mkdir(parents=True, exist_ok=True)
        with open(output_article, 'w', encoding='utf-8') as f:
            f.write(f"# {video_path.stem}\n\n")
            f.write("@全局:女声\n\n")
            f.write(full_text)
            f.write("\n")

        print(f"   ✅ 语音识别完成: {len(texts)} 段, {len(full_text)} 字")
        print(f"   📝 文章已保存: {output_article.name}")
        return True
    except Exception as e:
        print(f"   ❌ 语音识别失败: {e}")
        traceback.print_exc()
        return False


def build_fade_filter(duration: float, fade_duration: float) -> str:
    """构建淡入淡出滤镜"""
    if fade_duration <= 0 or duration <= fade_duration * 2:
        return ""
    fade_out_start = max(0, duration - fade_duration)
    return f",fade=t=in:st=0:d={fade_duration},fade=t=out:st={fade_out_start}:d={fade_duration}"


# 音色映射表
VOICE_MAP = {
    # 官方名称
    'Xiaoxiao': 'zh-CN-XiaoxiaoNeural',
    'Xiaoyi': 'zh-CN-XiaoyiNeural',
    'Yunxi': 'zh-CN-YunxiNeural',
    'Yunjian': 'zh-CN-YunjianNeural',
    'Yunxia': 'zh-CN-YunxiaNeural',
    'Yunyang': 'zh-CN-YunyangNeural',
    # 别名 - 女声
    '女声': 'zh-CN-XiaoxiaoNeural',
    '晓晓': 'zh-CN-XiaoxiaoNeural',
    '晓伊': 'zh-CN-XiaoyiNeural',
    '云夏': 'zh-CN-YunxiaNeural',
    '成熟女': 'zh-CN-XiaoyiNeural',
    '年轻女': 'zh-CN-YunxiaNeural',
    # 别名 - 男声
    '男声': 'zh-CN-YunyangNeural',
    '云扬': 'zh-CN-YunyangNeural',
    '云希': 'zh-CN-YunxiNeural',
    '云健': 'zh-CN-YunjianNeural',
    '成熟男': 'zh-CN-YunyangNeural',
    '年轻男': 'zh-CN-YunxiNeural',
    '新闻男': 'zh-CN-YunjianNeural',
}


def parse_voice_segments(text: str, default_voice: str = 'Xiaoxiao') -> list:
    """解析文章中的 @音色 标记

    返回: [(voice_key, content), ...]

    支持格式:
    - @女声: 内容
    - @男声: 内容
    - @晓晓: 内容
    - @云扬: 内容
    - @全局:女声  (设置默认音色，不生成音频)
    """
    segments = []
    current_voice = default_voice

    # 检查全局设置 @全局:音色
    global_match = re.search(r'@全局[:：](\w+)', text)
    if global_match:
        global_voice = global_match.group(1)
        if global_voice in VOICE_MAP:
            current_voice = global_voice
            print(f"   🎭 全局音色设置: {global_voice}")
        # 从文本中移除全局设置行，避免被当成普通标记处理
        text = re.sub(r'@全局[:：]\w+\n?', '', text)
    return segments, text, current_voice


# ── 插件系统 ──
class PluginManager:
    """简单的插件系统，扫描项目 plugins/ 目录下的 .py 文件"""

    def __init__(self, project_dir: Path):
        self.hooks = {
            'pre_parse_article': [],
            'post_parse_article': [],
            'pre_generate_audio': [],
            'post_generate_audio': [],
            'pre_generate_scene': [],
            'post_generate_scene': [],
            'pre_build_subtitle': [],
            'pre_merge': [],
        }
        self._load_plugins(project_dir)

    def _load_plugins(self, project_dir: Path):
        plugins_dir = project_dir / 'plugins'
        if not plugins_dir.exists():
            return

        import importlib.util
        for plugin_file in sorted(plugins_dir.glob('*.py')):
            if plugin_file.name.startswith('_'):
                continue
            try:
                spec = importlib.util.spec_from_file_location(
                    plugin_file.stem, str(plugin_file)
                )
                mod = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(mod)
                # 自动注册钩子
                for hook_name in self.hooks:
                    fn = getattr(mod, hook_name, None)
                    if callable(fn):
                        self.hooks[hook_name].append(fn)
                print(f"   🔌 插件加载: {plugin_file.name}")
            except Exception as e:
                print(f"   ⚠️  插件加载失败 {plugin_file.name}: {e}")

    def run(self, hook_name: str, *args, **kwargs):
        """执行某个钩子的所有插件函数"""
        results = []
        for fn in self.hooks.get(hook_name, []):
            try:
                result = fn(*args, **kwargs)
                if result is not None:
                    results.append(result)
            except Exception as e:
                print(f"   ⚠️  插件钩子 {hook_name} 执行失败: {e}")
        return results

    def run_first(self, hook_name: str, default, *args, **kwargs):
        """执行钩子，返回第一个非 None 结果，否则返回 default"""
        for fn in self.hooks.get(hook_name, []):
            try:
                result = fn(*args, **kwargs)
                if result is not None:
                    return result
            except Exception as e:
                print(f"   ⚠️  插件钩子 {hook_name} 执行失败: {e}")
        return default


def parse_article_segments(text: str, default_voice: str = 'Xiaoxiao') -> tuple:
    """解析文章，提取音色和图片分配信息

    返回: (segments, default_image)
    segments: [(voice_key, content, image_ref), ...]
    default_image: str 或 None (默认图片引用)

    支持格式:
    - @女声: 内容
    - @男声: 内容
    - @图:01 / @图片:02 / @img:03  - 指定图片
    - @默认图:01 / @默认图:bg.jpg   - 设置默认图片
    """
    segments = []
    current_voice = default_voice
    current_image = None  # 当前段落指定的图片
    default_image = None  # 全局默认图片

    # 检查全局设置
    # 1. 全局音色设置
    global_voice_match = re.search(r'@全局[:：](\w+)', text)
    if global_voice_match:
        global_voice = global_voice_match.group(1)
        if global_voice in VOICE_MAP:
            current_voice = global_voice
            default_voice = global_voice  # 同时更新 default_voice，空行重置时用
            print(f"   🎭 全局音色设置: {global_voice}")
        text = re.sub(r'@全局[:：]\w+\n?', '', text)

    # 2. 默认图片设置 @默认图:xx 或 @默认图片:xx
    default_img_match = re.search(r'@默认图(?:片)?[:：]\s*(\S+)', text)
    if default_img_match:
        default_image = default_img_match.group(1)
        print(f"   🖼️  默认图片设置: {default_image}")
        text = re.sub(r'@默认图(?:片)?[:：]\s*\S+\n?', '', text)

    # 初始化 current_image 为 default_image
    current_image = default_image

    # 按行分割并解析
    lines = text.split('\n')
    current_content = []

    for line in lines:
        line = line.strip()
        if not line:
            # 空行，保存当前段落
            if current_content:
                content = ' '.join(current_content).strip()
                if content and len(content) >= 3:
                    segments.append((current_voice, content, current_image))
                current_content = []
            # 重置当前图片和音色为默认值
            current_image = default_image
            current_voice = default_voice
            continue

        # 跳过 Markdown 标题行
        if line.startswith('#'):
            continue

        # 检查是否是图片标记行 @图:xx / @图片:xx / @img:xx
        img_match = re.match(r'@(图|图片|img)[:：]\s*(\S+)(.*)', line, re.IGNORECASE)
        if img_match:
            img_ref = img_match.group(2).strip()
            remaining_content = img_match.group(3).strip()

            # 保存之前的内容（使用之前的图片设置）
            if current_content:
                content = ' '.join(current_content).strip()
                if content and len(content) >= 3:
                    segments.append((current_voice, content, current_image))
                current_content = []

            # 更新当前图片设置
            current_image = img_ref

            # 如果有剩余内容，添加到当前段落
            if remaining_content and len(remaining_content) >= 3:
                current_content.append(remaining_content)
            continue

        # 检查是否是音色标记行 @音色:内容
        voice_match = re.match(r'@(\w+)[:：](.*)', line)
        if voice_match:
            voice_key = voice_match.group(1)
            content = voice_match.group(2).strip()

            # 如果是有效音色
            if voice_key in VOICE_MAP:
                # 先保存之前的内容
                if current_content:
                    prev_content = ' '.join(current_content).strip()
                    if prev_content and len(prev_content) >= 3:
                        segments.append((current_voice, prev_content, current_image))
                    current_content = []

                # 更新音色并添加当前内容
                current_voice = voice_key
                if content and len(content) >= 3:
                    # 检查内容中是否嵌套图片标记 @图:xx
                    img_in_content = re.match(r'@(图|图片|img)[:：]\s*(\S+)(.*)', content, re.IGNORECASE)
                    if img_in_content:
                        img_ref = img_in_content.group(2).strip()
                        remaining = img_in_content.group(3).strip()
                        current_image = img_ref  # 更新当前图片
                        if remaining and len(remaining) >= 3:
                            segments.append((current_voice, remaining, current_image))
                    else:
                        segments.append((current_voice, content, current_image))
            else:
                # 不是有效音色，当作普通内容
                line = re.sub(r'^#+\s*', '', line)
                line = re.sub(r'^\s*[-*+\d]\s+', '', line)
                if line:
                    current_content.append(line)
        else:
            # 普通内容行
            line = re.sub(r'^#+\s*', '', line)
            line = re.sub(r'^\s*[-*+\d]\s+', '', line)
            if line:
                current_content.append(line)

    # 保存最后一段
    if current_content:
        content = ' '.join(current_content).strip()
        if content and len(content) >= 3:
            segments.append((current_voice, content, current_image))

    return segments, default_image


async def generate_tts_with_retry(content: str, voice_id: str, output_path: Path, rate: str = '+0%', max_retries: int = 3) -> bool:
    """生成单段TTS音频，支持网络重试"""
    for attempt in range(1, max_retries + 1):
        try:
            communicate = edge_tts.Communicate(content, voice=voice_id, rate=rate)
            await communicate.save(str(output_path))
            return True
        except Exception as e:
            if attempt < max_retries:
                wait = attempt * 2  # 递增等待: 2s, 4s, 6s
                print(f"   ⚠️  TTS生成失败 (尝试 {attempt}/{max_retries}): {e}")
                print(f"   ⏳  {wait}秒后重试...")
                await asyncio.sleep(wait)
            else:
                print(f"   ❌ TTS生成最终失败 (已重试 {max_retries} 次): {e}")
                traceback.print_exc()
                return False
    return False


async def generate_audio_from_article(article_path: Path, output_dir: Path, voice: str = 'Xiaoxiao', rate: str = '+0%', video_mode: bool = False) -> tuple:
    """从文章自动生成音频，支持 @音色 多角色标记和 @图 图片分配

    Args:
        video_mode: True=生成单个合并音频(视频模式), False=生成分段音频(图片模式)

    Returns:
        (success, segments_info)
        success: bool - 是否成功
        segments_info: list of dict - 每段的音色和图片分配信息
    """
    if not EDGE_TTS_AVAILABLE:
        print("   ⚠️  未安装 edge_tts，跳过音频生成")
        return False, []

    # 读取文章
    with open(article_path, 'r', encoding='utf-8') as f:
        raw_text = f.read()

    # 统一换行符（处理 Windows CRLF）
    raw_text = raw_text.replace('\r\n', '\n').replace('\r', '\n')

    # 清理文本（保留 @标记）
    text = re.sub(r'```[\s\S]*?```', '', raw_text)
    text = re.sub(r'`[^`]*`', '', text)
    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'\n{3,}', '\n\n', text)

    if not text.strip():
        return False, []

    print(f"   🎙️  正在生成音频: {article_path.name}")

    # 解析带音色和图片标记的段落
    segments, default_image = parse_article_segments(text, default_voice=voice)

    if not segments:
        print("   ❌ 未找到有效内容")
        return False, []

    # 显示音色分配
    voice_summary = {}
    image_summary = {}
    for v, _, img in segments:
        voice_summary[v] = voice_summary.get(v, 0) + 1
        img_key = img if img else '顺序分配'
        image_summary[img_key] = image_summary.get(img_key, 0) + 1
    print(f"   🎭 音色分配: {', '.join([f'{v}({n}段)' for v, n in voice_summary.items()])}")
    if len(image_summary) > 1 or default_image:
        print(f"   🖼️  图片分配: {', '.join([f'{k}({n}段)' for k, n in image_summary.items()])}")

    try:
        segments_info = []
        semaphore = asyncio.Semaphore(3)

        async def _gen_one(i, voice_key, content, img_ref, output_path, label='场景'):
            """单个音频生成，受 semaphore 限流"""
            async with semaphore:
                voice_id = VOICE_MAP.get(voice_key, VOICE_MAP.get(voice, 'zh-CN-XiaoxiaoNeural'))
                success = await generate_tts_with_retry(content, voice_id, output_path, rate)
                if success:
                    segments_info.append({'voice': voice_key, 'image': img_ref or default_image, 'text': content, 'index': i})
                    img_info = f" [图:{img_ref}]" if img_ref else ""
                    print(f"      ✓ {label} {i}: [{voice_key}]{img_info} {content[:25]}...")
                return success

        if video_mode:
            # 视频模式：合并所有段落，但按音色分段生成后合并
            print(f"   📝 视频模式：合并 {len(segments)} 个音色段落")

            if len(segments) == 1:
                # 只有一个段落，直接生成
                voice_key, content, img_ref = segments[0]
                voice_id = VOICE_MAP.get(voice_key, VOICE_MAP.get(voice, 'zh-CN-XiaoxiaoNeural'))
                output_path = output_dir / "scene_01.mp3"
                if await generate_tts_with_retry(content, voice_id, output_path, rate):
                    print(f"   ✅ 音频生成完成: {output_path.name} ({len(content)} 字, {voice_key})")
                    segments_info.append({'voice': voice_key, 'image': img_ref or default_image, 'text': content, 'index': 1})
                else:
                    return False, []
            else:
                # 多音色，并行生成临时文件
                temp_files = []
                tasks = []
                for i, (voice_key, content, img_ref) in enumerate(segments, 1):
                    temp_path = output_dir / f"_temp_{i:02d}.mp3"
                    temp_files.append(temp_path)
                    tasks.append(_gen_one(i, voice_key, content, img_ref, temp_path, label='段落'))

                results = await asyncio.gather(*tasks)
                if not all(results):
                    # 清理已生成的临时文件
                    for tf in temp_files:
                        if tf.exists():
                            tf.unlink()
                    return False, []

                # 合并音频
                output_path = output_dir / "scene_01.mp3"

                # 验证临时文件
                missing = [str(tf) for tf in temp_files if not tf.exists() or tf.stat().st_size == 0]
                if missing:
                    print(f"   ❌ 临时音频文件缺失或为空: {missing}")
                    return False, []

                # 方法1: concat demuxer（不重新编码，速度最快）
                with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
                    for temp_file in temp_files:
                        f.write(f"file '{temp_file}'\n")
                    concat_file = f.name

                success = False
                try:
                    cmd = ['ffmpeg', '-y', '-f', 'concat', '-safe', '0',
                           '-i', concat_file, '-c', 'copy', str(output_path)]
                    result = subprocess.run(cmd, capture_output=True, text=True)
                    if result.returncode == 0:
                        success = True
                    else:
                        stderr = result.stderr.replace('\n', ' ')[:300]
                        print(f"   ⚠️  concat demuxer 失败: {stderr}")
                finally:
                    try:
                        os.unlink(concat_file)
                    except OSError:
                        pass

                # 方法2: concat filter（重新编码，更兼容）
                if not success:
                    print(f"   🔄 回退到 concat filter 合并...")
                    try:
                        inputs = []
                        filters = []
                        for idx, tf in enumerate(temp_files):
                            inputs.extend(['-i', str(tf)])
                            filters.append(f"[{idx}:a]")
                        filters_str = ''.join(filters) + f"concat=n={len(temp_files)}:v=0:a=1[a]"
                        cmd = ['ffmpeg', '-y'] + inputs + [
                            '-filter_complex', filters_str,
                            '-map', '[a]', '-c:a', 'libmp3lame', '-b:a', '192k',
                            str(output_path)
                        ]
                        result = subprocess.run(cmd, capture_output=True, text=True)
                        if result.returncode == 0:
                            success = True
                        else:
                            stderr = result.stderr.replace('\n', ' ')[:300]
                            print(f"   ⚠️  concat filter 也失败: {stderr}")
                    except Exception as e:
                        print(f"   ⚠️  concat filter 异常: {e}")

                # 清理临时文件
                for temp_file in temp_files:
                    if temp_file.exists():
                        temp_file.unlink()

                if success:
                    print(f"   ✅ 音频生成完成: {output_path.name} ({len(segments)} 个音色段落)")
                else:
                    print(f"   ❌ 音频合并失败，05_audio/ 为空")
                segments_info.sort(key=lambda x: x['index'])
                return success, segments_info

            segments_info.sort(key=lambda x: x['index'])
            return True, segments_info
        else:
            # 图片模式：并行生成每个段落音频
            print(f"   🚀 并行生成 {len(segments)} 段音频（最多3并发）...")
            tasks = []
            for i, (voice_key, content, img_ref) in enumerate(segments, 1):
                output_path = output_dir / f"scene_{i:02d}.mp3"
                tasks.append(_gen_one(i, voice_key, content, img_ref, output_path))

            results = await asyncio.gather(*tasks)
            if not all(results):
                return False, []

            # 并发 append 顺序不保证，按 index 排序后再返回
            segments_info.sort(key=lambda x: x['index'])
            print(f"   ✅ 音频生成完成: {len(segments)} 个场景")
            return True, segments_info
    except Exception as e:
        print(f"   ❌ 音频生成失败: {e}")
        traceback.print_exc()
        return False, []


def generate_publish_copy(project_dir: Path, api_key: str = None, base_url: str = None,
                          model: str = None, provider: str = 'kimi') -> bool:
    """根据文章内容自动生成多平台发布文案

    支持 provider: kimi / deepseek
    环境变量: KIMI_API_KEY / DEEPSEEK_API_KEY / OPENAI_API_KEY / LLM_BASE_URL / LLM_MODEL
    """
    import requests

    provider = (provider or 'kimi').lower()
    defaults = {
        'kimi': {'base_url': 'https://api.moonshot.cn/v1', 'model': 'moonshot-v1-8k', 'env_key': 'KIMI_API_KEY'},
        'deepseek': {'base_url': 'https://api.deepseek.com/v1', 'model': 'deepseek-chat', 'env_key': 'DEEPSEEK_API_KEY'}
    }
    cfg = defaults.get(provider, defaults['kimi'])

    article_dir = project_dir / '01_article'
    if not article_dir.exists():
        print("❌ 未找到 01_article/ 目录")
        return False

    article_path = _get_latest_article(article_dir)
    if not article_path:
        print("❌ 未找到文章文件")
        return False

    try:
        with open(article_path, 'r', encoding='utf-8') as f:
            article_text = f.read()
    except Exception as e:
        print(f"❌ 读取文章失败: {e}")
        return False

    clean_text = re.sub(r'@[^:\n]+[:：]', '', article_text)
    clean_text = re.sub(r'[#*`\n\r]', ' ', clean_text)
    clean_text = re.sub(r'\s+', ' ', clean_text).strip()[:1500]

    api_key = (api_key or os.environ.get(cfg['env_key']) or os.environ.get('OPENAI_API_KEY'))
    base_url = base_url or os.environ.get('LLM_BASE_URL', cfg['base_url'])
    model = model or os.environ.get('LLM_MODEL', cfg['model'])

    if not api_key:
        print(f"\n⚠️  未配置 {provider.upper()} API 密钥")
        print(f"   请设置环境变量: export {cfg['env_key']}='sk-...'")
        print("   或: export OPENAI_API_KEY='sk-...'")
        return False

    prompt = f"""你是一位资深短视频运营专家。请根据以下视频文案内容，生成适合不同平台的发布文案。

【视频文案】
{clean_text}

请按以下格式输出（不要输出其他内容）：

【抖音】
标题：（15-25字，抓眼球，带悬念或数字）
话题：（5-8个相关话题标签，格式 #标签）
简介：（50-80字，引导点赞关注）

【视频号】
标题：（15-25字，正能量或实用价值）
话题：（5-8个相关话题标签）
简介：（50-80字，适合朋友圈传播）

【B站】
标题：（20-30字，可带【】分类标签）
话题：（5-8个相关话题标签）
简介：（80-120字，可引导三连）
"""

    print("\n✍️  正在生成发布文案...")
    try:
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        payload = {
            'model': model,
            'messages': [{'role': 'user', 'content': prompt}],
            'temperature': 0.8,
            'max_tokens': 1500
        }
        response = requests.post(f"{base_url.rstrip('/')}/chat/completions", headers=headers, json=payload, timeout=60)
        response.raise_for_status()
        data = response.json()
        copy_text = data['choices'][0]['message']['content']
    except Exception as e:
        print(f"❌ API 调用失败: {e}")
        return False

    # 保存文案
    copy_path = article_dir / 'copy.md'
    with open(copy_path, 'w', encoding='utf-8') as f:
        f.write(f"# 发布文案\n\n")
        f.write(f"> 自动生成于 {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n")
        f.write(copy_text)
        f.write("\n")

    print(f"✅ 发布文案已保存: {copy_path.name}")
    # 打印摘要
    lines = copy_text.split('\n')
    for line in lines[:15]:
        if line.strip():
            print(f"   {line.strip()}")
    if len(lines) > 15:
        print(f"   ... ({len(lines)} 行)")
    return True


def _get_latest_article(article_dir: Path):
    """获取 01_article 目录下最新的文章文件（按修改时间）"""
    if not article_dir.exists():
        return None
    files = list(article_dir.glob('*.md')) + list(article_dir.glob('*.txt'))
    if not files:
        return None
    return max(files, key=lambda f: f.stat().st_mtime)


def _needs_realtime_search(title: str) -> bool:
    """根据标题关键词判断是否需要联网搜索获取实时数据"""
    realtime_keywords = [
        '今天', '今日', '最新', '刚刚', '实时', '热点', '热搜', '突发',
        '紧急', '快讯', '头条', '新闻', '战报', '比分', '结果', '预测',
        '行情', '走势', '股价', '股市', '开盘', '收盘', '涨停', '跌',
        '发布', '发布会', '上市', '发售', '开卖', '预售', '首销',
        '地震', '台风', '暴雨', '疫情', '政策', '规定', '公告',
        '比赛', '赛事', '决赛', '夺冠', '获胜', '进球', '世界杯',
        '奥斯卡', '金马', '票房', '票房榜', '排行榜'
    ]
    t = title.lower()
    for kw in realtime_keywords:
        if kw in t:
            return True
    return False


def auto_generate_article_from_title(title: str, output_dir: Path, api_key: str = None,
                                      base_url: str = None, model: str = None,
                                      provider: str = 'kimi', search_web: bool = False) -> Optional[Path]:
    """根据标题直接调用 LLM 生成口播文章

    支持 provider:
      - kimi:     Moonshot API (默认)
      - deepseek: DeepSeek API

    环境变量（优先级低于 CLI 参数）:
      - KIMI_API_KEY / DEEPSEEK_API_KEY / OPENAI_API_KEY
      - LLM_BASE_URL
      - LLM_MODEL

    返回: 生成的文章路径，失败返回 None
    """
    import requests

    # provider 默认配置
    provider = (provider or 'kimi').lower()
    defaults = {
        'kimi': {
            'base_url': 'https://api.moonshot.cn/v1',
            'model': 'moonshot-v1-8k',
            'env_key': 'KIMI_API_KEY'
        },
        'deepseek': {
            'base_url': 'https://api.deepseek.com/v1',
            'model': 'deepseek-chat',
            'env_key': 'DEEPSEEK_API_KEY'
        }
    }

    if provider not in defaults:
        print(f"⚠️  不支持的 provider: {provider}，可用: {', '.join(defaults.keys())}")
        return None

    cfg = defaults[provider]

    # 解析 API 配置：CLI 参数 > 专属环境变量 > 通用环境变量 > 默认值
    api_key = (api_key
               or os.environ.get(cfg['env_key'])
               or os.environ.get('OPENAI_API_KEY'))
    base_url = base_url or os.environ.get('LLM_BASE_URL', cfg['base_url'])
    model = model or os.environ.get('LLM_MODEL', cfg['model'])

    if not api_key:
        print(f"\n⚠️  未配置 {provider.upper()} API 密钥")
        print(f"   请设置环境变量: export {cfg['env_key']}='sk-...'")
        print("   或: export OPENAI_API_KEY='sk-...'")
        print(f"   或命令行指定: --llm-api-key 'sk-...'")
        return None

    print(f"\n✍️  正在使用 {provider.upper()} ({model}) 生成文章...")
    if search_web:
        print(f"   🌐 已开启联网搜索，获取实时数据...")

    prompt = f"""你是一位资深短视频文案专家。请根据以下标题，撰写一篇适合 AI 配音的口播文章。

【标题】{title}

要求：
1. 用中文撰写，口语化、有节奏感，适合朗读
2. 总字数控制在 300-600 字（约 1-2 分钟朗读时长）
3. 结构：开头吸引注意力 → 正文展开 → 结尾总结/引发思考
4. 每段不要太长，用空行分段
5. 不要出现 "大家好"、"我是 XX" 等固定开场白，直接进入内容
6. 适当使用 "你知道吗"、"说实话"、"其实" 等口语表达增加亲和力
7. 结尾不要添加 "关注我"、"点赞" 等引导语
8. 基于你对该话题的知识，写出有信息量的内容，不要编造具体数据

请直接输出文章正文，不需要输出标题和任何其他说明。"""

    try:
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        payload = {
            'model': model,
            'messages': [{'role': 'user', 'content': prompt}],
            'temperature': 0.7,
            'max_tokens': 2000
        }
        # 联网搜索：Kimi 和 DeepSeek 格式不同
        if search_web:
            if provider == 'kimi':
                payload['tools'] = [{'type': 'web_search', 'web_search': {'enable': True}}]
            elif provider == 'deepseek':
                payload['tools'] = [{'type': 'web_search'}]
        response = requests.post(
            f"{base_url.rstrip('/')}/chat/completions",
            headers=headers, json=payload, timeout=120
        )
        response.raise_for_status()
        data = response.json()
        article_text = data['choices'][0]['message']['content'].strip()
    except requests.exceptions.HTTPError as e:
        err_body = e.response.text[:200] if e.response else '未知'
        print(f"❌ API 错误: {e.response.status_code if e.response else '?'}")
        print(f"   详情: {err_body}")
        return None
    except Exception as e:
        print(f"❌ API 调用失败: {e}")
        return None

    # 保存文章（时间戳命名，避免覆盖）
    article_dir = output_dir / '01_article'
    article_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    article_path = article_dir / f'文章_{timestamp}.md'

    full_article = f"# {title}\n\n@全局:女声\n\n{article_text}\n"

    with open(article_path, 'w', encoding='utf-8') as f:
        f.write(full_article)

    print(f"✅ 文章已保存: {article_path}")
    preview_lines = [l.strip() for l in article_text.split('\n') if l.strip()][:3]
    for line in preview_lines:
        print(f"   {line[:60]}...")
    return article_path


def _extract_keywords_simple(text: str) -> str:
    """零依赖关键词提取：去掉口语词和停用词，提炼核心内容"""
    # 优先提取引号、书名号内的内容（通常是核心概念）
    quotes = re.findall(r'[《"「『]([^》"」』]+)[》」』"]', text)
    if quotes:
        best = max(quotes, key=len)
        clean = re.sub(r'[^一-鿿A-Za-z0-9\s]', '', best).strip()
        if len(clean) >= 2:
            return clean[:20]

    # 常见开头口语词
    prefixes = ['你知道吗', '说实话', '其实', '事实上', '简单来说',
                '值得注意的是', '有趣的是', '总的来说', '首先', '其次', '最后',
                '令人惊讶的是', '更重要的是', '除此之外', '另一方面']

    result = text
    for prefix in prefixes:
        if result.startswith(prefix):
            result = result[len(prefix):]
            break

    # 去掉标点，只保留中文、英文、数字、空格
    result = re.sub(r'[^一-鿿A-Za-z0-9\s]', '', result)

    # 去掉常见停用词（整词匹配）
    stop_phrases = ['一个', '一下', '一些', '一直', '非常', '可以', '进行',
                    '表示', '认为', '通过', '根据', '关于', '对于', '我们',
                    '知道', '看到', '听到', '想到', '觉得', '感觉', '发现',
                    '开始', '结束', '完成', '准备', '继续', '需要', '想要',
                    '应该', '必须', '可能', '也许', '大概', '差不多', '一起',
                    '还是', '只是', '已经', '但是', '因为', '所以', '如果',
                    '虽然', '不过', '并且', '或者', '以及', '还有', '就是',
                    '这样', '那样', '这里', '那里', '什么', '怎么', '为什么',
                    '多少', '多久', '哪里', '谁', '哪儿', '如何', '怎样',
                    '莫过于', '一帆风顺', '普通用户', '近年来', '最近一段时间',
                    '不得不说', '毫无疑问', '众所周知', '简单来说']
    for phrase in stop_phrases:
        result = result.replace(phrase, ' ')

    # 去掉单字停用词
    # 先无条件去掉最常见的"的"、"了"，然后对其他停用字做保守处理
    result = result.replace('的', ' ').replace('了', ' ')
    stop_chars = set('是在和与或就都而及对能会要去来到上下中被把给让从为以等着过将还只最更太真呢吧吗哦嗯哈呀哪')
    chars = list(result)
    filtered = []
    for i, c in enumerate(chars):
        if c in stop_chars:
            # 检查前后是否都是中文字（避免拆开复合词中间的字）
            prev_is_cn = i > 0 and '一' <= chars[i-1] <= '鿿'
            next_is_cn = i < len(chars)-1 and '一' <= chars[i+1] <= '鿿'
            if prev_is_cn and next_is_cn:
                # 在复合词中间，保留
                filtered.append(c)
            else:
                filtered.append(' ')
        else:
            filtered.append(c)
    result = ''.join(filtered)

    # 清理空格
    result = re.sub(r'\s+', ' ', result).strip()

    return result[:20] if result else text[:10]


def _extract_image_keywords(segments: list, article_text: str,
                            api_key: str = None, base_url: str = None,
                            model: str = None, provider: str = 'kimi') -> list:
    """为每个段落提取英文图片搜索关键词

    返回: List[str]，长度与 segments 相同
    """
    import requests

    provider = (provider or 'kimi').lower()
    defaults = {
        'kimi': {'base_url': 'https://api.moonshot.cn/v1', 'model': 'moonshot-v1-8k', 'env_key': 'KIMI_API_KEY'},
        'deepseek': {'base_url': 'https://api.deepseek.com/v1', 'model': 'deepseek-chat', 'env_key': 'DEEPSEEK_API_KEY'}
    }
    cfg = defaults.get(provider, defaults['kimi'])

    api_key = (api_key or os.environ.get(cfg['env_key']) or os.environ.get('OPENAI_API_KEY'))
    base_url = base_url or os.environ.get('LLM_BASE_URL', cfg['base_url'])
    model = model or os.environ.get('LLM_MODEL', cfg['model'])

    if not api_key:
        return []

    # 构建段落摘要（前80字）
    segment_summaries = []
    for i, (voice, content, _) in enumerate(segments):
        summary = content[:80].replace('\n', ' ')
        segment_summaries.append(f"段落{i+1}: {summary}")

    segments_text = "\n".join(segment_summaries)
    prompt = f"""你是一位视觉内容策划专家。请为以下口播文章的每个段落，提取一个最适合用于图片搜索的英文关键词。

要求：
1. 关键词必须具体、视觉化，能直接用于图库搜索
2. 每个段落只输出1个关键词，不要解释
3. 关键词尽量用英文，如果内容极具中国特色可用拼音
4. 输出格式严格如下（每行一个）：
   段落1: keyword1
   段落2: keyword2

文章内容段落：
{segments_text}

请严格按格式输出。"""

    try:
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        payload = {
            'model': model,
            'messages': [{'role': 'user', 'content': prompt}],
            'temperature': 0.3,
            'max_tokens': 1000
        }
        response = requests.post(
            f"{base_url.rstrip('/')}/chat/completions",
            headers=headers, json=payload, timeout=60
        )
        response.raise_for_status()
        data = response.json()
        result_text = data['choices'][0]['message']['content'].strip()
    except Exception as e:
        print(f"   ⚠️  关键词提取失败: {e}")
        return []

    # 解析返回的关键词
    keywords = []
    for line in result_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        # 匹配 "段落N: keyword" 或 "N. keyword" 或 "- keyword"
        match = re.match(r'(?:段落\d+[:：]|\d+[.．]|[-*])\s*(.+)', line)
        if match:
            kw = match.group(1).strip()
            # 清理可能的引号
            kw = kw.strip('"\'').strip()
            if kw:
                keywords.append(kw)

    # 如果解析出的关键词数量不足，用占位符补齐
    while len(keywords) < len(segments):
        idx = len(keywords)
        # 提取前20个非空白字符作为 fallback
        content = segments[idx][1].strip()
        fallback = re.sub(r'\s+', '_', content[:20]).strip('_')
        keywords.append(fallback or f"segment_{idx}")

    return keywords[:len(segments)]


def _download_image(keyword: str, provider: str, api_key: str, save_path: Path) -> bool:
    """下载图片到指定路径

    支持 provider:
      - pollinations: Pollinations.ai（免费，无需 key）
      - unsplash: Unsplash API（需 Access Key）
      - pexels: Pexels API（需 API Key）
    """
    import requests
    import urllib.parse

    provider = (provider or 'pollinations').lower()

    try:
        if provider == 'pollinations':
            # Pollinations.ai: 直接生成式图片，URL 编码 prompt
            import hashlib
            safe_kw = urllib.parse.quote(keyword)
            stable_seed = int(hashlib.md5(keyword.encode('utf-8')).hexdigest(), 16) % 10000
            url = f"https://image.pollinations.ai/prompt/{safe_kw}?width=1920&height=1080&nologo=true&seed={stable_seed}&enhance=true"
            response = requests.get(url, timeout=120)
            response.raise_for_status()
            content_type = response.headers.get('Content-Type', '')
            if not content_type.startswith('image/'):
                print(f"   ⚠️  Pollinations 返回非图片内容: {content_type[:50]}")
                return False
            save_path.write_bytes(response.content)
            return True

        elif provider == 'unsplash':
            if not api_key:
                print("   ⚠️  Unsplash 需要 API Key")
                return False
            search_url = f"https://api.unsplash.com/search/photos?query={urllib.parse.quote(keyword)}&per_page=1&orientation=landscape"
            headers = {'Authorization': f'Client-ID {api_key}'}
            search_resp = requests.get(search_url, headers=headers, timeout=30)
            search_resp.raise_for_status()
            data = search_resp.json()
            if not data.get('results'):
                print(f"   ⚠️  Unsplash 未找到图片: {keyword}")
                return False
            img_url = data['results'][0]['urls']['regular']
            img_resp = requests.get(img_url, timeout=60)
            img_resp.raise_for_status()
            content_type = img_resp.headers.get('Content-Type', '')
            if not content_type.startswith('image/'):
                print(f"   ⚠️  Unsplash 返回非图片内容: {content_type[:50]}")
                return False
            save_path.write_bytes(img_resp.content)
            return True

        elif provider == 'pexels':
            if not api_key:
                print("   ⚠️  Pexels 需要 API Key")
                return False
            search_url = f"https://api.pexels.com/v1/search?query={urllib.parse.quote(keyword)}&per_page=1&orientation=landscape"
            headers = {'Authorization': api_key}
            search_resp = requests.get(search_url, headers=headers, timeout=30)
            search_resp.raise_for_status()
            data = search_resp.json()
            if not data.get('photos'):
                print(f"   ⚠️  Pexels 未找到图片: {keyword}")
                return False
            img_url = data['photos'][0]['src']['large']
            img_resp = requests.get(img_url, timeout=60)
            img_resp.raise_for_status()
            content_type = img_resp.headers.get('Content-Type', '')
            if not content_type.startswith('image/'):
                print(f"   ⚠️  Pexels 返回非图片内容: {content_type[:50]}")
                return False
            save_path.write_bytes(img_resp.content)
            return True

        else:
            print(f"   ⚠️  不支持的图片 provider: {provider}")
            return False

    except requests.exceptions.HTTPError as e:
        err = e.response.text[:100] if e.response else str(e)
        print(f"   ⚠️  下载图片 HTTP 错误: {err}")
        return False
    except Exception as e:
        print(f"   ⚠️  下载图片失败: {e}")
        return False


def auto_generate_images_for_project(project_dir: Path,
                                      image_provider: str = 'pollinations',
                                      image_api_key: str = None,
                                      llm_provider: str = 'kimi',
                                      llm_api_key: str = None,
                                      llm_base_url: str = None,
                                      llm_model: str = None) -> int:
    """为项目文章自动配图

    返回: 成功插入的图片数量
    """
    article_dir = project_dir / '01_article'
    article_path = _get_latest_article(article_dir)
    if not article_path:
        print("❌ 未找到文章文件")
        return 0

    try:
        with open(article_path, 'r', encoding='utf-8') as f:
            article_text = f.read()
    except Exception as e:
        print(f"❌ 读取文章失败: {e}")
        return 0

    # 解析段落
    segments, default_image = parse_article_segments(article_text)
    if not segments:
        print("❌ 文章未解析出有效段落")
        return 0

    print(f"\n🎨 自动配图: 发现 {len(segments)} 个段落")
    print(f"   图片源: {image_provider}")

    # 提取关键词
    print("   🔍 正在提取图片关键词...")
    keywords = _extract_image_keywords(
        segments, article_text,
        api_key=llm_api_key, base_url=llm_base_url,
        model=llm_model, provider=llm_provider
    )
    if not keywords:
        print("   ℹ️  未配置 LLM API，直接用段落内容搜索图片（效果可能略差）")

    # 创建图片目录
    images_dir = project_dir / '03_images'
    images_dir.mkdir(exist_ok=True)

    # 下载图片
    downloaded_map = {}  # segment_idx -> filename
    for i, segment in enumerate(segments):
        if len(segment) != 3:
            print(f"   ⚠️  段落 {i} 格式异常，跳过")
            continue
        voice, content, _ = segment
        if i < len(keywords) and keywords[i]:
            keyword = keywords[i]
        else:
            # 无 LLM 关键词时，用规则提取核心关键词
            keyword = _extract_keywords_simple(content)
            if not keyword:
                keyword = f"segment_{i}"
        # 文件名与搜索词分开：文件名只取前15字，去掉中文标点，更简洁
        safe_kw = re.sub(r'[\\/:*?"<>|，。！？、；：""''（）【】]+', '', keyword).strip()[:15]
        if not safe_kw:
            safe_kw = f"segment_{i:02d}"
        filename = f"segment_{i+1:02d}_{safe_kw}.jpg"
        save_path = images_dir / filename

        print(f"   📥 [{i+1}/{len(segments)}] {keyword} → {filename}")
        success = _download_image(keyword, image_provider, image_api_key, save_path)
        if success:
            downloaded_map[i] = filename
            print(f"      ✅ 已下载")
        else:
            print(f"      ❌ 下载失败，跳过")
        # 避免触发限流
        if i < len(segments) - 1:
            time.sleep(0.5)

    if not downloaded_map:
        print("⚠️  没有成功下载任何图片")
        return 0

    # 在原文章中插入 @图: 标记
    print(f"\n📝 正在插入图片标记到文章...")
    # 注意：@图: 必须插入在 @音色: 之前，与 parse_article_segments 的解析顺序一致

    # 读取原文章并按行处理
    lines = article_text.split('\n')
    new_lines = []
    segment_idx = 0
    in_segment = False

    for line in lines:
        stripped = line.strip()

        # 跳过全局设置、默认图片、标题等元数据行
        if re.match(r'^#', stripped):
            new_lines.append(line)
            continue
        if re.match(r'^@(?:全局|默认图)', stripped):
            new_lines.append(line)
            continue
        if re.match(r'^@(图|图片|img)[:：]', stripped):
            # 跳过已有的图片标记，后面会重新插入
            continue

        # 检测段落边界（音色标记或空行）
        voice_match = re.match(r'^@(\w+)[:：]', stripped)
        if voice_match:
            voice_key = voice_match.group(1)
            if voice_key in VOICE_MAP:
                # 新段落开始，先插入图片标记
                if segment_idx in downloaded_map:
                    new_lines.append(f"@图:{downloaded_map[segment_idx]}")
                segment_idx += 1
                in_segment = True
                new_lines.append(line)
                continue

        if not stripped:
            # 空行表示段落结束
            in_segment = False
            new_lines.append(line)
            continue

        new_lines.append(line)

    # 处理最后一段（如果没有以空行结尾）
    # 这里 segment_idx 已经统计了所有音色标记，如果有遗漏在末尾补

    new_article = '\n'.join(new_lines)

    # 备份原文章
    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    backup_path = article_dir / f"文章_原稿_{timestamp}.md"
    shutil.copy(str(article_path), str(backup_path))
    print(f"   📋 原稿已备份: {backup_path.name}")

    # 保存修改后的文章
    with open(article_path, 'w', encoding='utf-8') as f:
        f.write(new_article)
    print(f"   ✅ 文章已更新: {article_path.name}")

    # 验证：重新解析，确保标记插入正确
    try:
        verify_segments, _ = parse_article_segments(new_article)
        img_markers = sum(1 for _, _, img in verify_segments if img)
        print(f"   🔍 验证: {len(verify_segments)} 段, {img_markers} 张图片分配")
    except Exception as e:
        print(f"   ⚠️  文章解析验证失败: {e}")

    print(f"\n{'='*60}")
    print(f"🎉 自动配图完成！共插入 {len(downloaded_map)} 张图片")
    print(f"{'='*60}")
    return len(downloaded_map)


def normalize_audio_loudness(input_path: Path, output_path: Path, target_lufs: float = -14.0) -> bool:
    """音频响度标准化（YouTube 标准 -14 LUFS）"""
    cmd = [
        'ffmpeg', '-y',
        '-i', str(input_path),
        '-af', f'loudnorm=I={target_lufs}:TP=-1.5:LRA=11',
        '-c:a', 'aac', '-b:a', '192k',
        str(output_path)
    ]
    try:
        result = run_ffmpeg(cmd, max_retries=1, check_output=False)
        if result.returncode == 0 and output_path.exists():
            return True
        print(f"   ⚠️  响度标准化失败，使用原音频")
        return False
    except Exception as e:
        print(f"   ⚠️  响度标准化异常: {e}")
        return False


def auto_generate_audio(project_dir: Path, voice: str = 'Xiaoxiao', rate: str = '+0%', force: bool = False, normalize: bool = False) -> tuple:
    """检查并自动生成音频

    Args:
        force: True=强制重新生成，删除已有音频

    Returns:
        (success, segments_info)
        success: bool - 是否成功
        segments_info: list of dict - 每段的音色和图片分配信息（图片模式）
    """
    audio_dir = project_dir / '05_audio'
    article_dir = project_dir / '01_article'

    # 查找文章文件
    latest_article = _get_latest_article(article_dir)

    # 解析文章段落数，用于判断是否需要重新生成
    article_segment_count = 0
    if latest_article:
        try:
            with open(latest_article, 'r', encoding='utf-8') as f:
                raw_text = f.read()
            raw_text = raw_text.replace('\r\n', '\n').replace('\r', '\n')
            text = re.sub(r'```[\s\S]*?```', '', raw_text)
            text = re.sub(r'`[^`]*`', '', text)
            text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
            text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
            text = re.sub(r'<[^>]+>', '', text)
            text = re.sub(r'\n{3,}', '\n\n', text)
            parsed_segments, _ = parse_article_segments(text, default_voice=voice)
            article_segment_count = len(parsed_segments)
        except Exception:
            pass

    existing_mp3s = sorted(audio_dir.glob('*.mp3')) if audio_dir.exists() else []

    # 如果已经有音频且非强制模式，检查段落数是否匹配
    if not force and existing_mp3s:
        if article_segment_count > 0 and len(existing_mp3s) != article_segment_count:
            print(f"🔄 文章段落数变化 ({len(existing_mp3s)} → {article_segment_count})，自动重新生成音频")
            for f in existing_mp3s:
                f.unlink()
        else:
            return False, []

    # 强制模式：删除已有音频
    if force and existing_mp3s:
        for f in existing_mp3s:
            f.unlink()
        print(f"🔄 强制重新生成音频，已清理旧音频")

    if not latest_article:
        return False, []

    # 检测模式：如果有视频目录且包含视频文件，则为视频模式
    videos_dir = project_dir / '04_videos'
    video_mode = videos_dir.exists() and any(videos_dir.iterdir())

    # 创建音频目录
    audio_dir.mkdir(exist_ok=True)

    # 选择最新的文章文件
    article_path = latest_article
    print(f"\n📄 发现文章: {article_path.name}")
    if video_mode:
        print("🎬 视频模式：生成单个完整音频")
    else:
        print("🖼️ 图片模式：生成多段音频")
    print("🎙️  正在自动生成音频...")

    # 运行异步生成
    try:
        success, segments_info = asyncio.run(generate_audio_from_article(article_path, audio_dir, voice, rate, video_mode))
    except Exception as e:
        print(f"   ❌ 生成失败: {e}")
        traceback.print_exc()
        return False, []

    # 响度标准化
    if normalize and success:
        print(f"\n🔊 音频响度标准化 (-14 LUFS)...")
        audio_files = sorted(audio_dir.glob('*.mp3'))
        for audio_file in audio_files:
            temp_normalized = audio_file.parent / f"_{audio_file.name}"
            if normalize_audio_loudness(audio_file, temp_normalized):
                audio_file.unlink()
                temp_normalized.rename(audio_file)
                print(f"   ✅ {audio_file.name}")
            else:
                if temp_normalized.exists():
                    temp_normalized.unlink()

    return success, segments_info


def find_image_by_ref(project_dir: Path, image_ref: str) -> Optional[Path]:
    """根据图片引用查找图片文件

    image_ref 可以是:
    - 数字: "01", "1" -> 查找 scene_01.jpg / 01.jpg / 1.jpg
    - 文件名: "bg.jpg" -> 查找 bg.jpg
    """
    if not image_ref:
        return None

    # 支持的图片目录和扩展名（包含大小写）
    img_dirs = ['03_images', '03_manual_images', '01_api_images']
    exts = ['.jpg', '.jpeg', '.png', '.webp', '.heic', '.gif',
            '.JPG', '.JPEG', '.PNG', '.WEBP', '.HEIC', '.GIF']

    for img_dir in img_dirs:
        img_dir_path = project_dir / img_dir
        if not img_dir_path.exists():
            continue

        # 如果是纯数字，尝试多种格式
        if image_ref.isdigit():
            num = int(image_ref)
            for ext in exts:
                # scene_XX.ext (scene_01.jpg)
                candidate = img_dir_path / f"scene_{num:02d}{ext}"
                if candidate.exists():
                    return candidate
                # XX.ext (01.jpg)
                candidate = img_dir_path / f"{num:02d}{ext}"
                if candidate.exists():
                    return candidate
                # X.ext (1.jpg) - 无前导零
                candidate = img_dir_path / f"{num}{ext}"
                if candidate.exists():
                    return candidate
        else:
            # 当作文件名处理
            for ext in exts:
                candidate = img_dir_path / f"{image_ref}{ext}"
                if candidate.exists():
                    return candidate
            # 如果用户已经带了扩展名
            candidate = img_dir_path / image_ref
            if candidate.exists():
                return candidate

    return None


def find_scenes(project_dir: Path, image_assignments: list = None) -> List[Scene]:
    """
    自动发现项目中的所有场景
    优先级：04_videos > 05_audio + 图片

    Args:
        image_assignments: 可选，每段音频对应的图片分配信息
            [{'index': 1, 'voice': '女声', 'image': '01'}, ...]
    """
    scenes = []

    # 先查找 04_videos 目录
    videos_dir = project_dir / '04_videos'
    if videos_dir.exists():
        video_files = sorted([f for f in videos_dir.iterdir()
                            if f.suffix.lower() in ['.mp4', '.mov', '.avi', '.mkv']])

        for i, video_file in enumerate(video_files, 1):
            duration = get_media_duration(str(video_file))
            # 尝试找对应的音频
            audio_path = project_dir / '05_audio' / f"scene_{i:02d}.mp3"
            if not audio_path.exists():
                audio_path = None

            # 从图片分配信息中获取字幕文本
            subtitle_text = ''
            if image_assignments:
                if i <= len(image_assignments):
                    subtitle_text = image_assignments[i - 1].get('text', '')
                # 视频数量少于 assignments 时，将剩余字幕合并到最后一个视频
                if len(video_files) < len(image_assignments) and i == len(video_files):
                    remaining_texts = [a.get('text', '') for a in image_assignments[i:]]
                    remaining = ' '.join([t for t in remaining_texts if t])
                    if remaining:
                        subtitle_text = (subtitle_text + ' ' + remaining).strip()

            scenes.append(Scene(
                index=i,
                audio_path=audio_path,
                video_path=video_file,
                image_path=None,
                subtitle=subtitle_text,
                duration=duration
            ))

    # 如果没有视频，查找音频+图片
    if not scenes:
        audio_dir = project_dir / '05_audio'
        if audio_dir.exists():
            audio_files = sorted([f for f in audio_dir.iterdir()
                                if f.suffix == '.mp3'])

            for i, audio_file in enumerate(audio_files, 1):
                duration = get_media_duration(str(audio_file))

                # 查找对应图片
                image_path = None

                # 优先使用图片分配信息
                if image_assignments and i <= len(image_assignments):
                    assignment = image_assignments[i - 1]  # index 从 1 开始
                    img_ref = assignment.get('image')
                    if img_ref:
                        image_path = find_image_by_ref(project_dir, img_ref)
                        if image_path:
                            print(f"   🖼️  场景 {i}: 使用指定图片 [{img_ref}] -> {image_path.name}")

                # 如果没有指定图片或找不到，使用顺序分配
                if not image_path:
                    exts = ['.jpg', '.jpeg', '.png', '.webp', '.heic', '.gif',
                            '.JPG', '.JPEG', '.PNG', '.WEBP', '.HEIC', '.GIF']
                    for img_dir in ['03_images', '03_manual_images', '01_api_images']:
                        img_dir_path = project_dir / img_dir
                        if img_dir_path.exists():
                            # 1. 尝试标准命名 scene_XX.ext
                            for ext in exts:
                                candidate = img_dir_path / f"scene_{i:02d}{ext}"
                                if candidate.exists():
                                    image_path = candidate
                                    break
                            # 2. 尝试数字命名 XX.ext
                            if not image_path:
                                for ext in exts:
                                    candidate = img_dir_path / f"{i:02d}{ext}"
                                    if candidate.exists():
                                        image_path = candidate
                                        break
                            # 3. 尝试无前导零 X.ext
                            if not image_path:
                                for ext in exts:
                                    candidate = img_dir_path / f"{i}{ext}"
                                    if candidate.exists():
                                        image_path = candidate
                                        break
                            # 4. 回退到目录中第 i 个图片（支持中文名任意命名）
                            if not image_path:
                                all_images = sorted([
                                    f for f in img_dir_path.iterdir()
                                    if f.is_file() and f.suffix.lower() in
                                       ['.jpg', '.jpeg', '.png', '.webp', '.gif', '.heic', '.bmp', '.tiff']
                                ])
                                # 有默认图时用默认图，否则按顺序取第 i 个
                                # 注意：这里需要获取 default_image 的值
                                # 但由于当前上下文没有 default_image，我们在外层处理
                                if i <= len(all_images):
                                    image_path = all_images[i - 1]
                            if image_path:
                                break

                # 从图片分配信息中获取字幕文本
                subtitle_text = ''
                if image_assignments and i <= len(image_assignments):
                    subtitle_text = image_assignments[i - 1].get('text', '')

                scenes.append(Scene(
                    index=i,
                    audio_path=audio_file,
                    video_path=None,
                    image_path=image_path,
                    subtitle=subtitle_text,
                    duration=duration
                ))

    return scenes


def create_scene_with_effects(
    scene: Scene,
    output_path: Path,
    resolution: Tuple[int, int],
    fps: int,
    add_subtitle: bool = False,
    subtitle_style: Dict = None,
    preview: bool = False,
    scene_fade: float = 0.0,
    subtitle_animation: str = 'none',
    subtitle_mode: str = 'sentence',
    subtitle_gap: float = 0.1
) -> bool:
    """创建单个场景视频，支持缩放效果、字幕、淡入淡出"""

    width, height = resolution

    if scene.video_path:
        # 视频素材 - 静音原视频，使用生成的配音音频
        video_input = str(scene.video_path)

        # 构建视频滤镜（缩放+字幕+淡入淡出）
        vf_filter = f"scale={width}:{height}:force_original_aspect_ratio=decrease,pad={width}:{height}:(ow-iw)/2:(oh-ih)/2:black"

        if add_subtitle and subtitle_style and scene.subtitle:
            if subtitle_mode == 'sentence' and scene.audio_path:
                audio_dur = get_media_duration(str(scene.audio_path))
                vf_filter += "," + build_sentence_subtitle_filter(
                    scene.subtitle, subtitle_style, width, height, audio_dur,
                    animation=subtitle_animation, subtitle_gap=subtitle_gap
                )
            else:
                vf_filter += "," + build_subtitle_filter(
                    scene.subtitle, subtitle_style, width, height,
                    animation=subtitle_animation
                )

        vf_filter += build_fade_filter(scene.duration, scene_fade)

        # 如果有生成的音频，使用生成的音频；否则保留原视频音频
        if scene.audio_path:
            # 获取音频时长
            audio_duration = get_media_duration(str(scene.audio_path))
            video_duration = scene.duration

            if audio_duration > video_duration:
                # 音频比视频长，需要循环视频
                loop_count = int(audio_duration / video_duration) + 1
                print(f"   🔄 视频循环: {loop_count} 次 (音频{audio_duration:.1f}s > 视频{video_duration:.1f}s)")

                cmd = [
                    'ffmpeg', '-y',
                    '-stream_loop', str(loop_count),
                    '-i', video_input,
                    '-i', str(scene.audio_path),
                    '-vf', vf_filter,
                    '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
                    '-pix_fmt', 'yuv420p',
                    '-c:a', 'aac', '-b:a', '192k',
                    '-map', '0:v:0', '-map', '1:a:0',
                    '-t', str(audio_duration),
                    '-shortest',
                    str(output_path)
                ]
            else:
                # 音频比视频短或相等，直接复制
                cmd = [
                    'ffmpeg', '-y',
                    '-i', video_input,
                    '-i', str(scene.audio_path),
                    '-vf', vf_filter,
                    '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
                    '-pix_fmt', 'yuv420p',
                    '-c:a', 'aac', '-b:a', '192k',
                    '-map', '0:v:0', '-map', '1:a:0',
                    '-t', str(audio_duration),
                    str(output_path)
                ]
        else:
            # 没有生成音频，直接复制原视频
            cmd = [
                'ffmpeg', '-y',
                '-i', video_input,
                '-vf', vf_filter,
                '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
                '-pix_fmt', 'yuv420p',
                '-c:a', 'copy',
                '-t', str(scene.duration),
                str(output_path)
            ]

    elif scene.image_path:
        # 图片素材 - 添加缩放动画效果
        duration = scene.duration
        total_frames = int(duration * fps)

        # 智能裁剪：等比放大到完全覆盖目标尺寸，居中裁剪，减少黑边
        crop_vf = f"scale={width}:{height}:force_original_aspect_ratio=increase,crop={width}:{height}:(iw-{width})/2:(ih-{height})/2,"

        # Ken Burns 效果：缓慢缩放和平移
        vf_filter = (
            crop_vf +
            f"zoompan=z='min(zoom+0.0005,1.1)':d={total_frames}:"
            f"x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)':s={width}x{height},"
            f"fps={fps},format=yuv420p,trim=duration={duration}"
        )

        # 添加字幕
        if add_subtitle and subtitle_style and scene.subtitle:
            if subtitle_mode == 'sentence' and scene.audio_path:
                audio_dur = get_media_duration(str(scene.audio_path))
                vf_filter += "," + build_sentence_subtitle_filter(
                    scene.subtitle, subtitle_style, width, height, audio_dur,
                    animation=subtitle_animation, subtitle_gap=subtitle_gap
                )
            else:
                vf_filter += "," + build_subtitle_filter(
                    scene.subtitle, subtitle_style, width, height,
                    animation=subtitle_animation
                )

        # 淡入淡出
        vf_filter += build_fade_filter(duration, scene_fade)

        # 编码参数：预览模式快速编码，正式模式高质量
        if preview:
            video_preset = 'ultrafast'
            video_crf = '28'
            audio_bitrate = '128k'
        else:
            video_preset = 'veryslow'
            video_crf = '15'
            audio_bitrate = '256k'

        # 添加音频（如果有）
        if scene.audio_path:
            cmd = [
                'ffmpeg', '-y',
                '-loop', '1',
                '-i', str(scene.image_path),
                '-i', str(scene.audio_path),
                '-vf', vf_filter,
                '-c:v', 'libx264',
                '-preset', video_preset,
                '-crf', video_crf,
                '-tune', 'stillimage',
                '-pix_fmt', 'yuv420p',
                '-c:a', 'aac',
                '-b:a', audio_bitrate,
                '-movflags', '+faststart',
                '-shortest',
                str(output_path)
            ]
        else:
            cmd = [
                'ffmpeg', '-y',
                '-loop', '1',
                '-i', str(scene.image_path),
                '-vf', vf_filter,
                '-c:v', 'libx264',
                '-preset', video_preset,
                '-crf', video_crf,
                '-tune', 'stillimage',
                '-pix_fmt', 'yuv420p',
                '-movflags', '+faststart',
                '-t', str(duration),
                '-an',
                str(output_path)
            ]
    else:
        return False

    try:
        result = run_ffmpeg(cmd, max_retries=2)
        if result.returncode != 0:
            print(f"   ⚠️ ffmpeg: {result.stderr[:200]}")
        return output_path.exists()
    except Exception as e:
        print(f"   ❌ 生成失败: {e}")
        traceback.print_exc()
        return False


def _generate_scene_worker(task):
    """并行场景生成工作函数"""
    scene, scene_output, width, height, fps, add_subtitle, subtitle_style, preview, scene_fade, subtitle_animation, subtitle_mode, subtitle_gap = task
    try:
        success = create_scene_with_effects(
            scene, scene_output, (width, height), fps,
            add_subtitle, subtitle_style, preview=preview,
            scene_fade=scene_fade, subtitle_animation=subtitle_animation,
            subtitle_mode=subtitle_mode, subtitle_gap=subtitle_gap
        )
        media_name = None
        if scene.image_path:
            media_name = scene.image_path.name
        elif scene.video_path:
            media_name = scene.video_path.name
        return {
            'index': scene.index,
            'success': success and scene_output.exists(),
            'output': scene_output,
            'error': None,
            'duration': scene.duration,
            'media_name': media_name
        }
    except Exception as e:
        media_name = None
        if scene.image_path:
            media_name = scene.image_path.name
        elif scene.video_path:
            media_name = scene.video_path.name
        traceback.print_exc()
        return {
            'index': scene.index,
            'success': False,
            'output': scene_output,
            'error': str(e),
            'duration': scene.duration,
            'media_name': media_name
        }


def add_transition(
    video1: Path,
    video2: Path,
    output: Path,
    transition_type: str,
    duration: float = 0.5,
    sfx_path: Optional[Path] = None
) -> bool:
    """添加转场效果，可选转场音效"""

    if transition_type == 'none' or not TRANSITIONS.get(transition_type):
        # 无转场，直接拼接
        return simple_concat([video1, video2], output)

    # 使用 xfade 滤镜实现转场
    # 获取两段视频的时长
    dur1 = get_media_duration(str(video1))
    dur2 = get_media_duration(str(video2))

    # 转场参数
    transition_name = TRANSITIONS[transition_type]
    offset = dur1 - duration

    # 构建命令
    inputs = ['-i', str(video1), '-i', str(video2)]
    audio_filter = f"[0:a][1:a]acrossfade=d={duration}[a]"

    # 混入转场音效
    if sfx_path and sfx_path.exists():
        sfx_dur = get_media_duration(str(sfx_path))
        use_dur = min(sfx_dur, duration * 2)
        inputs.extend(['-i', str(sfx_path)])
        audio_filter = (
            f"[0:a][1:a]acrossfade=d={duration}[mix];"
            f"[2:a]atrim=0:{use_dur},volume=0.6,adelay=0s[se];"
            f"[mix][se]amix=inputs=2:duration=first[a]"
        )

    cmd = [
        'ffmpeg', '-y',
        *inputs,
        '-filter_complex',
        f"[0:v][1:v]xfade=transition={transition_name}:duration={duration}:offset={offset}[v];"
        + audio_filter,
        '-map', '[v]', '-map', '[a]',
        '-c:v', 'libx264', '-preset', 'veryslow', '-crf', '15',
        '-pix_fmt', 'yuv420p',
        '-c:a', 'aac', '-b:a', '256k',
        str(output)
    ]

    try:
        result = run_ffmpeg(cmd, max_retries=2)
        if result.returncode != 0:
            print(f"   ⚠️  转场失败，使用简单拼接: {result.stderr[:100]}")
            return simple_concat([video1, video2], output)
        return output.exists()
    except Exception as e:
        print(f"   ⚠️  转场失败: {e}，使用简单拼接")
        return simple_concat([video1, video2], output)


def simple_concat(videos: List[Path], output: Path) -> bool:
    """简单拼接视频，先尝试快速复制拼接，失败则回退到重新编码"""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
        for video in videos:
            # concat demuxer 要求对单引号进行转义
            escaped = str(video).replace("'", "'\\''")
            f.write(f"file '{escaped}'\n")
        concat_file = f.name

    try:
        # 先尝试 -c copy（最快，不重新编码）
        cmd = [
            'ffmpeg', '-y', '-f', 'concat', '-safe', '0',
            '-i', concat_file,
            '-c', 'copy',
            str(output)
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode == 0 and output.exists() and output.stat().st_size > 1024:
            return True

        # 失败则回退到重新编码（兼容不同分辨率/帧率/编码参数）
        print(f"   🔄 复制拼接失败，回退到重新编码...")
        cmd = [
            'ffmpeg', '-y', '-f', 'concat', '-safe', '0',
            '-i', concat_file,
            '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
            '-pix_fmt', 'yuv420p',
            '-c:a', 'aac', '-b:a', '192k',
            str(output)
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        return result.returncode == 0 and output.exists() and output.stat().st_size > 1024
    finally:
        os.unlink(concat_file)


def add_intro_outro(
    main_video: Path,
    intro: Optional[Path],
    outro: Optional[Path],
    output: Path,
    transition: str = 'fade'
) -> bool:
    """添加片头片尾"""

    videos = []

    if intro and intro.exists():
        videos.append(intro)

    videos.append(main_video)

    if outro and outro.exists():
        videos.append(outro)

    if len(videos) == 1:
        shutil.copy(str(main_video), str(output))
        return True

    # 使用转场连接
    if len(videos) == 2:
        return add_transition(videos[0], videos[1], output, transition)

    if len(videos) == 3:
        temp = Path(tempfile.mktemp(suffix='.mp4'))
        try:
            if add_transition(videos[0], videos[1], temp, transition):
                return add_transition(temp, videos[2], output, transition)
            return False
        finally:
            if temp.exists():
                temp.unlink()

    return False


def add_bgm(video: Path, bgm: Path, output: Path, volume: float = 0.3) -> bool:
    """添加背景音乐（智能裁剪循环 + 淡入淡出）

    - BGM 比视频短：循环播放，循环衔接处做淡入淡出，结尾淡出让位
    - BGM 比视频长：裁剪到视频时长，结尾淡出让位
    """

    video_dur = get_media_duration(str(video))
    bgm_dur = get_media_duration(str(bgm))

    # 淡入淡出时长
    fade_in = min(2.0, bgm_dur * 0.2)
    fade_out = min(3.0, video_dur * 0.15)

    if bgm_dur < video_dur:
        # BGM 比视频短：循环播放 + 淡入淡出衔接
        loops = int(video_dur / bgm_dur) + 1
        # 对 BGM 做淡入淡出处理，循环后在视频结尾处整体淡出
        bgm_filter = (
            f"[1:a]volume={volume},"
            f"afade=t=in:st=0:d={fade_in},"
            f"afade=t=out:st={max(0,bgm_dur-fade_in)}:d={fade_in},"
            f"aloop=loop={loops}:size={int(bgm_dur*48000)},"
            f"atrim=0:{video_dur},"
            f"afade=t=out:st={max(0,video_dur-fade_out)}:d={fade_out}[bgm]"
        )
        print(f"   🎵 BGM 循环: {bgm_dur:.1f}s → {video_dur:.1f}s ({loops}次, 淡入{fade_in:.1f}s/淡出{fade_out:.1f}s)")
    else:
        # BGM 比视频长：裁剪 + 淡出让位
        bgm_filter = (
            f"[1:a]volume={volume},"
            f"afade=t=in:st=0:d={fade_in},"
            f"atrim=0:{video_dur},"
            f"afade=t=out:st={max(0,video_dur-fade_out)}:d={fade_out}[bgm]"
        )
        print(f"   🎵 BGM 裁剪: {bgm_dur:.1f}s → {video_dur:.1f}s (淡入{fade_in:.1f}s/淡出{fade_out:.1f}s)")

    mix_filter = "[0:a][bgm]amix=inputs=2:duration=first[a]"

    cmd = [
        'ffmpeg', '-y',
        '-i', str(video),
        '-i', str(bgm),
        '-filter_complex',
        f"{bgm_filter};{mix_filter}",
        '-map', '0:v', '-map', '[a]',
        '-c:v', 'copy', '-c:a', 'aac', '-b:a', '192k',
        str(output)
    ]

    try:
        result = run_ffmpeg(cmd, max_retries=2)
        if result.returncode != 0:
            print(f"   ⚠️ BGM添加失败，使用原视频")
            shutil.copy(str(video), str(output))
            return True
        return output.exists()
    except Exception as e:
        print(f"   ⚠️ BGM添加失败: {e}")
        traceback.print_exc()
        shutil.copy(str(video), str(output))
        return True


def add_watermark(
    video: Path,
    watermark: Path,
    position: str,
    output: Path,
    opacity: float = 0.75
) -> bool:
    """添加水印/Logo"""

    positions = {
        'top-left': '10:10',
        'top-right': 'W-w-10:10',
        'bottom-left': '10:H-h-10',
        'bottom-right': 'W-w-10:H-h-10',
        'center': '(W-w)/2:(H-h)/2',
    }
    pos = positions.get(position, 'W-w-10:H-h-10')

    cmd = [
        'ffmpeg', '-y',
        '-i', str(video),
        '-i', str(watermark),
        '-filter_complex',
        f"[1:v]format=rgba,colorchannelmixer=aa={opacity}[wm];"
        f"[0:v][wm]overlay={pos}[v]",
        '-map', '[v]', '-map', '0:a',
        '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
        '-c:a', 'copy',
        str(output)
    ]

    try:
        result = run_ffmpeg(cmd, max_retries=2)
        if result.returncode != 0:
            print(f"   ⚠️ 水印添加失败: {result.stderr[:100]}")
            shutil.copy(str(video), str(output))
            return True
        return output.exists()
    except Exception as e:
        print(f"   ⚠️ 水印添加失败: {e}")
        traceback.print_exc()
        shutil.copy(str(video), str(output))
        return True


def generate_dual_version(source_video: Path, output_video: Path, target_resolution: str) -> bool:
    """生成相反比例版本（横屏↔竖屏），居中裁剪保留核心内容"""
    try:
        tw, th = map(int, target_resolution.split('x'))
    except ValueError:
        print(f"   ⚠️  无效目标分辨率: {target_resolution}")
        return False

    # ffmpeg: 等比放大让短边填满目标，然后居中裁剪
    vf = (
        f"scale=iw*max({tw}/iw\\,{th}/ih):ih*max({tw}/iw\\,{th}/ih),"
        f"crop={tw}:{th}"
    )

    cmd = [
        'ffmpeg', '-y', '-i', str(source_video),
        '-vf', vf,
        '-c:v', 'libx264', '-preset', 'fast', '-crf', '18',
        '-pix_fmt', 'yuv420p',
        '-c:a', 'copy',
        '-movflags', '+faststart',
        str(output_video)
    ]

    try:
        result = run_ffmpeg(cmd, max_retries=1, check_output=False)
        if result.returncode != 0:
            print(f"   ⚠️  双版本生成失败: {result.stderr[:200]}")
            return False
        if output_video.exists():
            size_mb = output_video.stat().st_size / (1024 * 1024)
            dur = get_media_duration(str(output_video))
            print(f"   ✅ 双版本完成: {output_video.name} ({dur:.1f}s, {size_mb:.1f}MB)")
            return True
        return False
    except Exception as e:
        print(f"   ⚠️  双版本异常: {e}")
        traceback.print_exc()
        return False


def pre_check_project(project_dir: Path, args) -> Tuple[bool, List[str]]:
    """项目素材预检，返回 (是否通过, 错误信息列表)

    检查项:
    - 项目目录结构
    - 文章文件存在性
    - 图片/视频素材存在性
    - 字体文件可用性
    - 水印文件存在性
    - 分辨率格式合法性
    - BGM文件存在性
    - 输出目录可写性
    """
    errors = []
    warnings = []

    # 1. 项目目录
    if not project_dir.exists():
        errors.append(f"项目目录不存在: {project_dir}")
        return False, errors

    # 先检测模式（供后续检查使用）
    images_dir = project_dir / '03_images'
    videos_dir = project_dir / '04_videos'
    img_exts = ['.jpg', '.jpeg', '.png', '.webp', '.gif', '.bmp', '.tiff']

    if videos_dir.exists() and any(f.suffix.lower() in ['.mp4', '.mov', '.avi', '.mkv'] for f in videos_dir.iterdir()):
        mode = 'video'
    elif images_dir.exists() and any(f.suffix.lower() in img_exts for f in images_dir.iterdir()):
        mode = 'image'
    else:
        errors.append("缺少素材: 03_images/ 无图片，且 04_videos/ 无视频")
        mode = None

    # 2. 文章文件
    article_dir = project_dir / '01_article'
    if not article_dir.exists():
        errors.append("缺少 01_article/ 目录（无文章素材）")
    else:
        latest_article = _get_latest_article(article_dir)
        if not latest_article:
            # 视频模式下允许无文章（可使用语音识别）
            if mode == 'video':
                warnings.append("未找到文章，视频模式可使用语音识别生成字幕")
            else:
                errors.append("01_article/ 目录下没有找到 .md 或 .txt 文章文件")

    # 3. 图片/视频素材（模式已在上面检测）

    # 4. 字体可用性（检查常用中文字体）
    font_candidates = [
        '/System/Library/Fonts/STHeiti Medium.ttc',
        '/System/Library/Fonts/PingFang.ttc',
        '/System/Library/Fonts/Hiragino Sans GB.ttc',
        '/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc',
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    ]
    font_found = any(Path(f).exists() for f in font_candidates)
    if not font_found:
        warnings.append("未检测到常用中文字体，字幕可能显示为方框")

    # 5. 水印文件
    watermark = getattr(args, 'watermark', None)
    if watermark:
        watermark_path = Path(watermark)
        if not watermark_path.is_absolute():
            watermark_path = project_dir / watermark_path
        if not watermark_path.exists():
            errors.append(f"水印文件不存在: {watermark_path}")

    # 6. 分辨率格式
    resolution = getattr(args, 'resolution', '')
    if resolution:
        try:
            w, h = map(int, resolution.split('x'))
            if w <= 0 or h <= 0 or w > 7680 or h > 4320:
                errors.append(f"分辨率数值异常: {resolution}")
        except ValueError:
            errors.append(f"分辨率格式错误: {resolution}，应为 宽x高 如 1920x1080")

    # 7. BGM
    bgm = getattr(args, 'bgm', None)
    if bgm:
        bgm_path = Path(bgm)
        if not bgm_path.exists():
            errors.append(f"指定BGM不存在: {bgm_path}")
    else:
        bgm_dir = project_dir / '02_bgm'
        if bgm_dir.exists():
            music_files = [f for f in bgm_dir.iterdir() if f.suffix.lower() in ['.mp3', '.wav', '.aac', '.flac', '.m4a', '.ogg']]
            if not music_files:
                warnings.append("02_bgm/ 目录存在但没有音乐文件")

    # 8. 输出目录可写
    final_dir = project_dir / '07_final'
    try:
        final_dir.mkdir(parents=True, exist_ok=True)
        test_file = final_dir / '._write_test'
        test_file.write_text('test')
        test_file.unlink()
    except Exception:
        errors.append(f"输出目录不可写: {final_dir}")

    # 打印报告
    print(f"\n🔍 素材预检报告: {project_dir.name}")
    print(f"{'─'*50}")
    if mode:
        print(f"   模式: {'🎬 视频模式' if mode == 'video' else '🖼️  图片模式'}")
    if latest_article:
        print(f"   文章: {latest_article.name}")
    if font_found:
        print(f"   字体: ✅ 可用")
    if not warnings and not errors:
        print(f"   状态: ✅ 全部通过")
    for w in warnings:
        print(f"   ⚠️  {w}")
    for e in errors:
        print(f"   ❌ {e}")
    print(f"{'─'*50}")

    return len(errors) == 0, errors


def compute_file_hash(path: Path) -> str:
    """计算文件MD5哈希，用于增量更新比较"""
    try:
        import hashlib
        h = hashlib.md5()
        with open(path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def load_build_manifest(scenes_dir: Path) -> Dict:
    """加载构建清单，记录每个场景的输入哈希和参数"""
    manifest_path = scenes_dir / '.build_manifest.json'
    if manifest_path.exists():
        try:
            with open(manifest_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            pass
    return {}


def save_build_manifest(scenes_dir: Path, manifest: Dict):
    """保存构建清单"""
    manifest_path = scenes_dir / '.build_manifest.json'
    try:
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


def should_rebuild_scene(
    scene: Scene,
    scene_output: Path,
    manifest: Dict,
    width: int, height: int, fps: int,
    add_subtitle: bool, subtitle_style_key: str,
    scene_fade: float, subtitle_animation: str, subtitle_mode: str
) -> bool:
    """判断场景是否需要重新生成（增量更新）"""
    # 如果输出文件不存在或损坏，需要重建
    if not scene_output.exists() or scene_output.stat().st_size <= 1024:
        return True

    scene_key = f"scene_{scene.index:02d}"
    if scene_key not in manifest:
        return True

    record = manifest[scene_key]

    # 检查输入文件哈希
    input_files = []
    if scene.image_path:
        input_files.append(('image', scene.image_path))
    if scene.video_path:
        input_files.append(('video', scene.video_path))
    if scene.audio_path:
        input_files.append(('audio', scene.audio_path))

    for key, path in input_files:
        current_hash = compute_file_hash(path)
        if record.get(key) != current_hash:
            return True

    # 检查生成参数
    params = {
        'width': width,
        'height': height,
        'fps': fps,
        'add_subtitle': add_subtitle,
        'subtitle_style': subtitle_style_key if subtitle_style_key else 'none',
        'scene_fade': scene_fade,
        'subtitle_animation': subtitle_animation,
        'subtitle_mode': subtitle_mode,
    }
    if record.get('params') != params:
        return True

    return False


def process_project(
    project_dir: Path,
    args
) -> Optional[Path]:
    """处理单个项目"""

    tee = None
    try:
        tee = setup_logging(project_dir, preview_mode=getattr(args, 'preview', False))

        # 预览模式
        preview_mode = getattr(args, 'preview', False)
        start_time = time.time()
        stage_times = {}

        # 素材预检
        if not getattr(args, 'skip_pre_check', False):
            passed, errors = pre_check_project(project_dir, args)
            if not passed:
                print(f"\n❌ 预检未通过，终止处理")
                return None

        print(f"\n{'='*60}")
        if preview_mode:
            print(f"👁️  预览模式: {project_dir.name}")
            print(f"   只生成第一个场景，快速预览效果")
        else:
            print(f"🎬 处理项目: {project_dir.name}")
        print(f"{'='*60}")

        # 加载插件
        plugin_mgr = PluginManager(project_dir)

        # 检查输出目录
        final_dir = project_dir / '07_final'
        final_dir.mkdir(exist_ok=True)

        # 创建场景片段目录（保存中间产物，支持断点续传）
        scenes_dir = project_dir / '06_scenes'
        scenes_dir.mkdir(exist_ok=True)

        if preview_mode:
            output_path = final_dir / 'preview.mp4'
        else:
            if args.output:
                output_name = args.output
            else:
                now = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
                output_name = f"{project_dir.name}_{now}.mp4"
            output_path = final_dir / output_name

        # 视频模式逻辑处理
        videos_dir = project_dir / '04_videos'
        is_video_mode = videos_dir.exists() and any(videos_dir.iterdir())
        keep_original_audio = False

        if is_video_mode and not preview_mode:
            article_dir = project_dir / '01_article'
            article_files = list(article_dir.glob('*.md')) + list(article_dir.glob('*.txt')) if article_dir.exists() else []
            video_files = sorted([f for f in videos_dir.iterdir()
                                if f.suffix.lower() in ['.mp4', '.mov', '.avi', '.mkv']])

            if not article_files:
                # 路径 A: 无文章，询问是否使用语音识别
                use_whisper = getattr(args, 'whisper_transcribe', False)
                if not use_whisper:
                    choice = input("\n🎙️ 未找到文章，是否使用本地语音识别生成字幕？ [Y/n]: ").strip().replace('\r', '').lower()
                    use_whisper = choice != 'n'

                if use_whisper and video_files:
                    article_path = article_dir / '文章.md'
                    if transcribe_video_with_whisper(video_files[0], article_path):
                        article_files = [article_path]
                        keep_original_audio = True
                        print("🎬 语音识别完成，保留视频原声")
                    else:
                        print("❌ 语音识别失败，无法继续")
                        return None
                elif not article_files:
                    print("❌ 视频模式需要文章或语音识别")
                    return None
            else:
                # 有文章：直接使用 AI 配音替换原声
                print("🎬 使用 AI 配音替换原声")

        # 自动从文章生成音频（如果没有音频但有文章）
        # 返回音频分段信息，包含音色和图片分配
        audio_t0 = time.time()
        regenerate_audio = getattr(args, 'regenerate_audio', False)

        if keep_original_audio:
            # 保留原声：只解析文章获取字幕信息，不生成音频
            latest_article = _get_latest_article(project_dir / '01_article')
            if latest_article:
                try:
                    with open(latest_article, 'r', encoding='utf-8') as f:
                        raw_text = f.read()
                    raw_text = raw_text.replace('\r\n', '\n').replace('\r', '\n')
                    text = re.sub(r'```[\s\S]*?```', '', raw_text)
                    text = re.sub(r'`[^`]*`', '', text)
                    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
                    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
                    text = re.sub(r'<[^>]+>', '', text)
                    text = re.sub(r'\n{3,}', '\n\n', text)
                    parsed_segments, _ = parse_article_segments(text, default_voice=getattr(args, 'voice', 'Xiaoxiao'))
                    segments_info = [
                        {'voice': v, 'image': img, 'text': content, 'index': i}
                        for i, (v, content, img) in enumerate(parsed_segments, 1)
                    ]
                    # 插件钩子：文章解析后
                    plugin_results = plugin_mgr.run('post_parse_article', segments_info)
                    if plugin_results:
                        segments_info = plugin_results[-1]
                except Exception:
                    pass
            audio_success = True
            stage_times['audio'] = 0
        else:
            audio_success, segments_info = auto_generate_audio(
                project_dir,
                voice=getattr(args, 'voice', 'Xiaoxiao'),
                rate=getattr(args, 'rate', '+0%'),
                force=regenerate_audio,
                normalize=getattr(args, 'normalize_audio', False)
            )
            stage_times['audio'] = time.time() - audio_t0

        # 如果已有音频但没有分段信息，从文章中解析字幕文本
        if not segments_info and not keep_original_audio:
            article_dir = project_dir / '01_article'
            latest_article = _get_latest_article(article_dir)
            if latest_article:
                try:
                    with open(latest_article, 'r', encoding='utf-8') as f:
                        raw_text = f.read()
                        raw_text = raw_text.replace('\r\n', '\n').replace('\r', '\n')
                        text = re.sub(r'```[\s\S]*?```', '', raw_text)
                        text = re.sub(r'`[^`]*`', '', text)
                        text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
                        text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
                        text = re.sub(r'<[^>]+>', '', text)
                        text = re.sub(r'\n{3,}', '\n\n', text)
                        parsed_segments, _ = parse_article_segments(text, default_voice=getattr(args, 'voice', 'Xiaoxiao'))
                        segments_info = [
                            {'voice': v, 'image': img, 'text': content, 'index': i}
                            for i, (v, content, img) in enumerate(parsed_segments, 1)
                        ]
                        # 插件钩子：文章解析后
                        plugin_results = plugin_mgr.run('post_parse_article', segments_info)
                        if plugin_results:
                            segments_info = plugin_results[-1]
                except Exception:
                    pass

        # 发现场景（传递图片分配信息）
        scenes = find_scenes(project_dir, image_assignments=segments_info)
        if not scenes:
            print("❌ 未找到任何场景素材")
            return None

        # 预览模式只取第一个场景
        if preview_mode:
            scenes = scenes[:1]
            print(f"📁 发现 {len(scenes)} 个场景 (预览模式只取第一个)")
        else:
            print(f"📁 发现 {len(scenes)} 个场景")
        print(f"📂 场景片段保存到: {scenes_dir}")

        # 字幕样式
        subtitle_style = None
        if args.subtitle:
            # 未显式指定时，根据分辨率自动适配
            if args.subtitle_style is None:
                try:
                    width, height = map(int, args.resolution.split('x'))
                    if height > width:
                        args.subtitle_style = 'tiktok'
                        print(f"   📝 竖屏自动适配字幕样式: tiktok")
                    else:
                        args.subtitle_style = 'news'
                        print(f"   📝 横屏自动适配字幕样式: news")
                except (ValueError, IndexError):
                    args.subtitle_style = 'news'
                    print(f"   ⚠️  分辨率格式异常，使用默认字幕样式: news")
            subtitle_style = SUBTITLE_STYLES.get(args.subtitle_style, SUBTITLE_STYLES['news'])
            print(f"📝 字幕样式: {args.subtitle_style}")

        if not preview_mode:
            print(f"✨ 转场效果: {args.transition}")

        # 处理每个场景（支持断点续传 + 增量更新）
        scene_videos = []
        failed_scenes = []
        skipped_scenes = []
        pending_tasks = []
        width, height = map(int, args.resolution.split('x'))

        # 预览模式：降低分辨率加速出片
        if preview_mode:
            orig_res = f"{width}x{height}"
            preview_width = 480
            scale = preview_width / width
            width = preview_width
            height = int(height * scale)
            # 确保高度为偶数（H.264 要求）
            if height % 2 != 0:
                height += 1
            print(f"   预览分辨率: {width}x{height} (原 {orig_res})")

        # 加载构建清单（增量更新）
        build_manifest = load_build_manifest(scenes_dir) if not preview_mode else {}
        new_manifest = {}

        for i, scene in enumerate(scenes):
            scene_output = scenes_dir / f"scene_{scene.index:02d}.mp4"

            # 预览模式跳过断点续传（每次预览都重新生成）
            if not preview_mode:
                # 断点续传：检查是否已存在且有效
                if scene_output.exists() and scene_output.stat().st_size > 1024:
                    # 验证文件是否有效（可播放）
                    try:
                        verify_duration = get_media_duration(str(scene_output))
                        if verify_duration > 0:
                            # 增量更新：进一步检查输入文件和参数是否变化
                            scene_fade = getattr(args, 'scene_fade', 0.0)
                            subtitle_animation = getattr(args, 'subtitle_animation', 'none')
                            subtitle_mode = getattr(args, 'subtitle_mode', 'sentence')
                            needs_rebuild = should_rebuild_scene(
                                scene, scene_output, build_manifest,
                                width, height, args.fps,
                                args.subtitle, args.subtitle_style if subtitle_style else None,
                                scene_fade, subtitle_animation, subtitle_mode
                            )
                            if not needs_rebuild:
                                scene_videos.append(scene_output)
                                skipped_scenes.append(scene.index)
                                print(f"\n[{i+1}/{len(scenes)}] 场景 {scene.index:02d}...")
                                print(f"   ⏭️  未变更，跳过生成 ({verify_duration:.1f}s)")
                                # 保留原manifest记录
                                scene_key = f"scene_{scene.index:02d}"
                                if scene_key in build_manifest:
                                    new_manifest[scene_key] = build_manifest[scene_key]
                                continue
                            else:
                                print(f"\n[{i+1}/{len(scenes)}] 场景 {scene.index:02d}...")
                                print(f"   🔄 素材或参数变更，重新生成")
                    except:
                        pass  # 文件损坏，重新生成

            print(f"\n[{i+1}/{len(scenes)}] 场景 {scene.index:02d}...")

            if scene.video_path:
                print(f"   🎥 视频: {scene.video_path.name}")
            elif scene.image_path:
                print(f"   🖼️  图片: {scene.image_path.name}")
            else:
                print(f"   ⚠️  跳过: 无素材")
                failed_scenes.append(scene.index)
                continue

            if scene.audio_path:
                print(f"   🎵 音频: {scene.audio_path.name}")

            scene_fade = getattr(args, 'scene_fade', 0.0)
            subtitle_animation = getattr(args, 'subtitle_animation', 'none')
            subtitle_mode = getattr(args, 'subtitle_mode', 'sentence')
            subtitle_gap = getattr(args, 'subtitle_gap', 0.1)
            pending_tasks.append((scene, scene_output, width, height, args.fps, args.subtitle, subtitle_style, preview_mode, scene_fade, subtitle_animation, subtitle_mode, subtitle_gap))

        # 执行生成（支持并行）
        total_pending = len(pending_tasks)
        scene_t0 = time.time()
        try:
            from tqdm import tqdm
        except ImportError:
            tqdm = None

        if not args.no_parallel and total_pending > 1 and not preview_mode:
            print(f"\n🚀 并行生成 {total_pending} 个场景（最多3并发）...")
            pbar = tqdm(total=total_pending, desc="场景生成", unit="个", file=sys.stdout) if tqdm else None
            _parallel_results = []
            with concurrent.futures.ProcessPoolExecutor(max_workers=3) as executor:
                futures = [executor.submit(_generate_scene_worker, task) for task in pending_tasks]
                for future in concurrent.futures.as_completed(futures):
                    result = future.result()
                    idx = result['index']
                    media = result.get('media_name') or '无素材'
                    if result['success']:
                        _parallel_results.append(result)
                        if pbar:
                            pbar.update(1)
                        else:
                            print(f"   ✅ [{len(_parallel_results)}/{total_pending}] 场景 {idx:02d} ({media}) 完成 ({result['duration']:.1f}s)")
                    else:
                        failed_scenes.append(idx)
                        if pbar:
                            pbar.update(1)
                        else:
                            print(f"   ❌ [{len(_parallel_results)+1}/{total_pending}] 场景 {idx:02d} ({media}) 失败: {result['error'] or '未知错误'}")
            if pbar:
                pbar.close()
            # 按场景 index 排序后再取输出，确保合并顺序正确
            _parallel_results.sort(key=lambda r: r['index'])
            scene_videos = [r['output'] for r in _parallel_results]
        else:
            # 顺序生成（预览模式也用顺序）
            pbar = tqdm(total=total_pending, desc="场景生成", unit="个", file=sys.stdout) if tqdm else None
            for task_idx, task in enumerate(pending_tasks, 1):
                scene, scene_output, width, height, fps, add_subtitle, subtitle_style, preview, scene_fade, subtitle_animation, subtitle_mode, subtitle_gap = task
                media = scene.image_path.name if scene.image_path else (scene.video_path.name if scene.video_path else '无素材')
                try:
                    success = create_scene_with_effects(
                        scene, scene_output, (width, height), fps,
                        add_subtitle, subtitle_style, preview=preview,
                        scene_fade=scene_fade, subtitle_animation=subtitle_animation,
                        subtitle_mode=subtitle_mode, subtitle_gap=subtitle_gap
                    )
                    if success and scene_output.exists():
                        scene_videos.append(scene_output)
                        if pbar:
                            pbar.update(1)
                        else:
                            print(f"   ✅ [{task_idx}/{total_pending}] 场景 {scene.index:02d} ({media}) 完成 ({scene.duration:.1f}s)")
                    else:
                        if pbar:
                            pbar.update(1)
                        else:
                            print(f"   ❌ [{task_idx}/{total_pending}] 场景 {scene.index:02d} ({media}) 生成失败")
                        failed_scenes.append(scene.index)
                except Exception as e:
                    if pbar:
                        pbar.update(1)
                    else:
                        print(f"   ❌ [{task_idx}/{total_pending}] 场景 {scene.index:02d} ({media}) 异常: {e}")
                    traceback.print_exc()
                    failed_scenes.append(scene.index)
            if pbar:
                pbar.close()
        stage_times['scenes'] = time.time() - scene_t0

        # 保存构建清单（增量更新）
        if not preview_mode:
            for scene in scenes:
                scene_output = scenes_dir / f"scene_{scene.index:02d}.mp4"
                scene_key = f"scene_{scene.index:02d}"
                if scene_output.exists() and scene_output.stat().st_size > 1024:
                    record = {'output_exists': True}
                    if scene.image_path:
                        record['image'] = compute_file_hash(scene.image_path)
                    if scene.video_path:
                        record['video'] = compute_file_hash(scene.video_path)
                    if scene.audio_path:
                        record['audio'] = compute_file_hash(scene.audio_path)
                    record['params'] = {
                        'width': width,
                        'height': height,
                        'fps': args.fps,
                        'add_subtitle': args.subtitle,
                        'subtitle_style': args.subtitle_style if subtitle_style else 'none',
                        'scene_fade': getattr(args, 'scene_fade', 0.0),
                        'subtitle_animation': getattr(args, 'subtitle_animation', 'none'),
                        'subtitle_mode': getattr(args, 'subtitle_mode', 'sentence'),
                    }
                    new_manifest[scene_key] = record
            save_build_manifest(scenes_dir, new_manifest)

        # 汇总生成结果
        print(f"\n{'─'*60}")
        print(f"📊 场景生成统计:")
        print(f"   ✅ 成功: {len(scene_videos)} 个")
        if skipped_scenes:
            print(f"   ⏭️  跳过(未变更): {len(skipped_scenes)} 个 {skipped_scenes}")
        if failed_scenes:
            print(f"   ❌ 失败: {len(failed_scenes)} 个 {failed_scenes}")
        print(f"{'─'*60}")

        if not scene_videos:
            print("❌ 没有成功生成任何场景")
            return None

        # 预览模式：直接复制第一个场景，跳过合并/转场/片头片尾/BGM
        if preview_mode:
            main_video = scene_videos[0]
            preview_duration = getattr(args, 'preview_duration', 5.0)
            video_dur = get_media_duration(str(main_video))

            if video_dur > preview_duration:
                # 裁剪到预览时长
                print(f"   ✂️  裁剪到 {preview_duration}s (原 {video_dur:.1f}s)")
                cmd = [
                    'ffmpeg', '-y',
                    '-i', str(main_video),
                    '-t', str(preview_duration),
                    '-c', 'copy',
                    str(output_path)
                ]
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode != 0:
                    shutil.copy(str(main_video), str(output_path))
            else:
                shutil.copy(str(main_video), str(output_path))

            final_duration = get_media_duration(str(output_path))
            final_size = output_path.stat().st_size / (1024 * 1024)

            print(f"\n{'='*60}")
            print("✅ 预览视频生成完成!")
            print(f"{'='*60}")
            print(f"📁 输出: {output_path}")
            print(f"⏱️  时长: {final_duration:.1f} 秒")
            print(f"📦 大小: {final_size:.1f} MB")
            print(f"🎞️  场景: 1 个（预览模式）")
            print(f"📂 场景片段: {scenes_dir}")

            return output_path

        # 查找转场音效
        sfx_path = None
        if getattr(args, 'sfx', False):
            sfx_dir = project_dir / '02_sfx'
            if sfx_dir.exists():
                sfx_exts = ['.mp3', '.wav', '.aac', '.m4a', '.ogg']
                sfx_files = sorted([f for f in sfx_dir.iterdir() if f.suffix.lower() in sfx_exts])
                if sfx_files:
                    sfx_path = sfx_files[0]
                    print(f"🔊 自动发现转场音效: {sfx_path.name}")

        # 合并场景（带转场）
        print(f"\n🎞️  合并场景（转场: {args.transition}）...")

        merge_state_path = final_dir / '.merge_state.json'
        merge_t0 = time.time()

        with tempfile.TemporaryDirectory() as temp_dir_str:
            temp_dir = Path(temp_dir_str)
            if len(scene_videos) == 1:
                main_video = scene_videos[0]
            elif args.transition == 'none' or not TRANSITIONS.get(args.transition):
                # 无转场，一次性拼接（速度最快）
                concat_output = temp_dir / 'concat_all.mp4'
                print(f"   🔄 [1/1] 拼接 {len(scene_videos)} 个场景（无需重新编码）...")
                if simple_concat(scene_videos, concat_output):
                    main_video = concat_output
                    print(f"   ✅ 拼接完成")
                else:
                    print(f"   ❌ 拼接失败")
                    return None
            else:
                # 有转场，逐步合并，支持断点续传
                resume_idx = 1
                main_video = scene_videos[0]

                # 检查是否有未完成的合并
                if merge_state_path.exists():
                    try:
                        with open(merge_state_path, 'r', encoding='utf-8') as f:
                            state = json.load(f)
                        state_scenes = state.get('scene_names', [])
                        current_scenes = [s.name for s in scene_videos]
                        if state_scenes == current_scenes and state.get('transition') == args.transition:
                            partial_file = final_dir / state.get('output', '')
                            if partial_file.exists():
                                duration = get_media_duration(str(partial_file))
                                if duration > 0:
                                    resume_idx = state.get('resume_idx', 1)
                                    main_video = partial_file
                                    print(f"⏭️  恢复未完成的合并，从 [{resume_idx}/{len(scene_videos)}] 继续")
                    except Exception:
                        traceback.print_exc()
                        pass  # 状态文件损坏，重新合并

                merge_pbar = tqdm(total=len(scene_videos)-1, desc="合并转场", unit="个", file=sys.stdout, initial=resume_idx-1) if tqdm else None
                for i in range(resume_idx, len(scene_videos)):
                    next_video = scene_videos[i]
                    partial_file = final_dir / f'_partial_merged_{i}.mp4'
                    from_name = main_video.name if hasattr(main_video, 'name') else str(main_video)
                    to_name = next_video.name

                    if merge_pbar:
                        merge_pbar.set_postfix({"from": from_name[:15], "to": to_name[:15]})
                    else:
                        print(f"   🔄 [{i+1}/{len(scene_videos)}] 合并: {from_name} → {to_name}")

                    prev_video = main_video
                    success = add_transition(prev_video, next_video, partial_file, args.transition, args.transition_duration, sfx_path)

                    if success:
                        if merge_pbar:
                            merge_pbar.update(1)
                        # 保存合并状态
                        state = {
                            'scene_names': [s.name for s in scene_videos],
                            'transition': args.transition,
                            'output': partial_file.name,
                            'resume_idx': i + 1,
                            'timestamp': datetime.datetime.now().isoformat()
                        }
                        with open(merge_state_path, 'w', encoding='utf-8') as f:
                            json.dump(state, f, ensure_ascii=False, indent=2)

                        # 清理旧的 partial 文件（如果是 partial 且不是当前输入）
                        if isinstance(prev_video, Path) and prev_video.name.startswith('_partial_merged_'):
                            try:
                                prev_video.unlink()
                            except OSError:
                                pass

                        main_video = partial_file
                    else:
                        if merge_pbar:
                            merge_pbar.close()
                        print(f"   ❌ 合并失败，终止")
                        return None
                if merge_pbar:
                    merge_pbar.close()

            stage_times['merge'] = time.time() - merge_t0

            # 添加片头片尾
            intro_path = Path(args.intro) if getattr(args, 'intro', None) else None
            outro_path = Path(args.outro) if getattr(args, 'outro', None) else None
            intro_text = getattr(args, 'intro_text', None)
            outro_text = getattr(args, 'outro_text', None)

            if intro_path or outro_path or intro_text or outro_text:
                print("🎬 添加片头片尾...")
                with_intro_outro = temp_dir / 'with_intro_outro.mp4'

                # 如果提供了文字，自动生成片头/片尾视频
                if intro_text:
                    intro_path = temp_dir / '_auto_intro.mp4'
                    cw, ch = map(int, args.resolution.split('x'))
                    generate_text_video(intro_text, intro_path, (cw, ch), duration=3.0, fps=args.fps)
                if outro_text:
                    outro_path = temp_dir / '_auto_outro.mp4'
                    cw, ch = map(int, args.resolution.split('x'))
                    generate_text_video(outro_text, outro_path, (cw, ch), duration=3.0, fps=args.fps)

                if add_intro_outro(main_video, intro_path, outro_path, with_intro_outro, args.transition):
                    main_video = with_intro_outro
                    print("✅ 片头片尾添加完成")

            # 添加背景音乐
            bgm_path = None
            if args.bgm:
                # 显式指定 BGM 路径
                bgm_path = Path(args.bgm)
            else:
                # 自动查找项目 02_bgm/ 目录下的音乐文件
                bgm_dir = project_dir / '02_bgm'
                if bgm_dir.exists():
                    music_exts = ['.mp3', '.wav', '.aac', '.flac', '.m4a', '.ogg']
                    music_files = sorted([f for f in bgm_dir.iterdir() if f.suffix.lower() in music_exts])
                    if music_files:
                        bgm_path = music_files[0]
                        print(f"🎵 自动发现 BGM: {bgm_path.name}")

            if bgm_path and bgm_path.exists():
                print(f"🎵 添加背景音乐 (音量: {int(args.bgm_volume * 100)}%)...")
                with_bgm = temp_dir / 'with_bgm.mp4'
                if add_bgm(main_video, bgm_path, with_bgm, args.bgm_volume):
                    main_video = with_bgm
                    print("✅ BGM添加完成")

            # 添加水印
            watermark_path = getattr(args, 'watermark', None)
            if watermark_path:
                watermark_path = Path(watermark_path)
                if not watermark_path.is_absolute():
                    watermark_path = project_dir / watermark_path
            if watermark_path and watermark_path.exists():
                print(f"🏷️ 添加水印 ({getattr(args, 'watermark_position', 'bottom-right')})...")
                watermarked = temp_dir / 'watermarked.mp4'
                if add_watermark(main_video, watermark_path, getattr(args, 'watermark_position', 'bottom-right'), watermarked):
                    main_video = watermarked
                    print("✅ 水印添加完成")

            # 复制到输出位置
            shutil.copy(str(main_video), str(output_path))

            # 双版本生成（横竖屏）
            if getattr(args, 'dual_version', False) and not preview_mode:
                cw, ch = map(int, args.resolution.split('x'))
                alt_res = f"{ch}x{cw}"
                if cw > ch:
                    alt_name = output_path.stem + '_portrait' + output_path.suffix
                    label = '竖屏版'
                else:
                    alt_name = output_path.stem + '_landscape' + output_path.suffix
                    label = '横屏版'
                alt_path = output_path.parent / alt_name
                print(f"\n📱 生成双版本 ({label} {alt_res})...")
                generate_dual_version(output_path, alt_path, alt_res)

            # 清理中间合并文件和状态
            for f in final_dir.glob('_partial_merged_*.mp4'):
                try:
                    f.unlink()
                except OSError:
                    pass
            if merge_state_path.exists():
                try:
                    merge_state_path.unlink()
                except OSError:
                    pass

            # 显示结果
            final_duration = get_media_duration(str(output_path))
            final_size = output_path.stat().st_size / (1024 * 1024)
            total_time = time.time() - start_time

            print(f"\n{'='*60}")
            print("✅ 视频合成完成!")
            print(f"{'='*60}")
            print(f"📁 输出: {output_path}")
            print(f"⏱️  时长: {final_duration:.1f} 秒")
            print(f"📦 大小: {final_size:.1f} MB")
            print(f"🎞️  场景: {len(scene_videos)}/{len(scenes)}")
            if skipped_scenes:
                print(f"⏭️  跳过: {len(skipped_scenes)} 个（已存在）")
            if failed_scenes:
                print(f"❌ 失败: {len(failed_scenes)} 个 {failed_scenes}")
            print(f"📂 场景片段: {scenes_dir}")

            # 详细报告
            print(f"\n{'─'*60}")
            print("📊 生成详细报告:")
            print(f"{'─'*60}")
            print(f"  ⏱️  总耗时:       {total_time:.1f}s")
            if 'audio' in stage_times:
                print(f"  🎙️  音频生成:     {stage_times['audio']:.1f}s")
            if 'scenes' in stage_times:
                print(f"  🎬 场景生成:     {stage_times['scenes']:.1f}s")
            if 'merge' in stage_times:
                print(f"  🔗 合并转场:     {stage_times['merge']:.1f}s")
            print(f"  📹 最终时长:     {final_duration:.1f}s")
            print(f"  📦 文件大小:     {final_size:.1f} MB")
            if final_duration > 0:
                print(f"  📊 平均码率:     {final_size * 8 / final_duration:.1f} Mbps")
            # 场景详情
            if len(scene_videos) > 1:
                print(f"\n  📋 各场景详情:")
                for i, sv in enumerate(scene_videos, 1):
                    sv_duration = get_media_duration(str(sv))
                    sv_size = sv.stat().st_size / (1024 * 1024)
                    print(f"     场景{i:02d}: {sv_duration:5.1f}s | {sv_size:5.1f}MB")
            print(f"{'─'*60}")

            # 自动提取封面
            if not preview_mode and output_path.exists():
                cover_path = final_dir / 'cover.jpg'
                try:
                    # 提取中间帧作为封面
                    mid_time = final_duration / 2 if final_duration > 0 else 1
                    cmd = [
                        'ffmpeg', '-y',
                        '-i', str(output_path),
                        '-ss', str(mid_time),
                        '-vframes', '1',
                        '-q:v', '2',
                        str(cover_path)
                    ]
                    result = run_ffmpeg(cmd, max_retries=2, check_output=False)
                    if cover_path.exists():
                        print(f"\n🖼️  封面已提取: {cover_path.name}")
                except Exception as e:
                    pass  # 封面提取失败不影响主流程

            return output_path

    finally:
        # 恢复 stdout 并关闭日志（临时目录由 TemporaryDirectory 自动清理）
        if tee:
            sys.stdout = tee.terminal
            tee.close()


def import_ppt_project(ppt_path: Path, project_dir: Path) -> bool:
    """导入PPT/Keynote项目：提取图片、备注文本，生成项目结构

    依赖:
    - python-pptx (pip install python-pptx)
    - soffice (LibreOffice) 可选，用于将PPT每页转为PNG截图

    流程:
    1. 提取PPT中所有嵌入图片 -> 03_images/
    2. 提取每页备注/文本 -> 01_article/文章.md
    3. 如有soffice，自动将每页转为PNG截图 -> 03_images/slide_XX.png
    4. 初始化项目配置
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ 缺少依赖: pip install python-pptx")
        return False

    if not ppt_path.exists():
        print(f"❌ PPT文件不存在: {ppt_path}")
        return False

    print(f"\n{'='*60}")
    print(f"📊 PPT导入: {ppt_path.name}")
    print(f"{'='*60}")

    try:
        prs = Presentation(str(ppt_path))
    except Exception as e:
        print(f"❌ 无法读取PPT: {e}")
        return False

    # 创建项目目录
    project_dir.mkdir(parents=True, exist_ok=True)
    (project_dir / '01_article').mkdir(exist_ok=True)
    (project_dir / '03_images').mkdir(exist_ok=True)
    (project_dir / '02_bgm').mkdir(exist_ok=True)

    slide_count = len(prs.slides)
    print(f"📑 共 {slide_count} 页幻灯片")

    # 提取嵌入图片
    img_dir = project_dir / '03_images'
    extracted_images = 0
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    ext = image.ext if image.ext else 'png'
                    img_path = img_dir / f"ppt_image_{i:02d}_{extracted_images+1:02d}.{ext}"
                    img_path.write_bytes(image.blob)
                    extracted_images += 1
                except Exception:
                    pass
    if extracted_images > 0:
        print(f"🖼️  提取嵌入图片: {extracted_images} 张")

    # 生成文章.md（从备注和文本）
    article_lines = [f"# {ppt_path.stem}\n"]
    has_notes = False
    for i, slide in enumerate(prs.slides, 1):
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_text_frame.text.strip()
        if notes_text:
            has_notes = True
            article_lines.append(f"\n@图:slide_{i:02d}\n")
            article_lines.append(f"{notes_text}\n")
        else:
            # 无备注时，提取该页主要文本
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    t = shape.text.strip()
                    # 过滤页码和过短文本
                    if len(t) > 3:
                        texts.append(t)
            if texts:
                article_lines.append(f"\n@图:slide_{i:02d}\n")
                article_lines.append(f"{'。'.join(texts[:3])}。\n")

    article_path = project_dir / '01_article' / '文章.md'
    with open(article_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(article_lines))
    print(f"📝 生成文章: {article_path.name} ({'含演讲备注' if has_notes else '从页面文本提取'})")

    # 尝试用soffice将每页转为PNG
    soffice_cmd = shutil.which('soffice') or shutil.which('libreoffice')
    if soffice_cmd:
        print(f"🔄 使用 {soffice_cmd} 转换幻灯片为图片...")
        import tempfile, subprocess
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            cmd = [
                soffice_cmd, '--headless', '--convert-to', 'png',
                '--outdir', str(tmp_path), str(ppt_path)
            ]
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode == 0:
                png_files = sorted(tmp_path.glob('*.png'))
                for idx, png in enumerate(png_files, 1):
                    dest = img_dir / f"slide_{idx:02d}.png"
                    shutil.copy(str(png), str(dest))
                print(f"✅ 幻灯片截图: {len(png_files)} 张 -> 03_images/")
            else:
                print(f"⚠️  soffice 转换失败: {result.stderr[:200]}")
    else:
        print(f"\n⚠️  未检测到 LibreOffice (soffice)")
        print(f"   如需PPT每页截图，请:")
        print(f"   1. 安装 LibreOffice: brew install --cask libreoffice")
        print(f"   2. 或手动用 PowerPoint/WPS 另存为图片到 03_images/")
        print(f"   当前已提取所有嵌入图片和文本，可直接运行生成")

    # 写入项目配置
    config = {
        "mode": "image",
        "resolution": "1920x1080",
        "fps": 30,
        "subtitle": True,
        "subtitle_style": "news",
        "subtitle_animation": "none",
        "transition": "fade",
        "voice": "Xiaoxiao",
        "transition_duration": 0.5,
        "rate": "+18%",
        "scene_fade": 0.0,
        "bgm_volume": 0.25,
        "watermark": None,
        "watermark_position": "bottom-right",
        "sfx": False,
        "dual_version": False,
        "created": datetime.datetime.now().isoformat()
    }
    config_path = project_dir / '.video_config.json'
    with open(config_path, 'w', encoding='utf-8') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)
    print(f"⚙️  初始化配置: {config_path.name}")

    print(f"\n{'='*60}")
    print(f"✅ PPT导入完成!")
    print(f"{'='*60}")
    print(f"📁 项目: {project_dir}")
    print(f"📝 文章: 01_article/文章.md")
    print(f"🖼️  图片: 03_images/")
    print(f"\n下一步:")
    print(f"  python3 01_核心脚本/video_generator_pro.py -p {project_dir}")
    return True


def init_project_wizard(project_dir: Path, template: str = None) -> bool:
    """交互式项目初始化向导

    Args:
        template: 可选，使用预设模板（news/food/tutorial/education）
    """
    # 检查是否使用模板
    template_config = None
    if template and template in PROJECT_TEMPLATES:
        template_config = PROJECT_TEMPLATES[template]
        print(f"\n{'='*60}")
        print(f"🚀 项目初始化向导 - 使用模板: {template}")
        print(f"{'='*60}")
    else:
        print(f"\n{'='*60}")
        print("🚀 项目初始化向导")
        print(f"{'='*60}")

    print(f"📁 项目路径: {project_dir}")

    if project_dir.exists() and any(project_dir.iterdir()):
        print(f"\n⚠️  目录已存在且不为空: {project_dir}")
        confirm = input("是否继续初始化? [y/N]: ").strip().replace('\r', '').lower()
        if confirm != 'y':
            print("已取消")
            return False

    project_dir.mkdir(parents=True, exist_ok=True)

    if template_config:
        # 使用模板预设
        mode = template_config['mode']
        resolution = template_config['resolution']
        fps = template_config['fps']
        subtitle_style = template_config['subtitle_style']
        transition = template_config['transition']
        voice = template_config['voice']
        print(f"\n📋 模板配置:")
        print(f"  模式: {mode}")
        print(f"  分辨率: {resolution}")
        print(f"  字幕: {subtitle_style}")
        print(f"  转场: {transition}")
        print(f"  音色: {voice}")
    else:
        # 1. 选择模式
        print("\n📋 请选择项目模式:")
        print("  1) 🖼️  图片模式 - 文章 + 图片素材")
        print("  2) 🎥 视频模式 - 文章 + 原视频素材")
        print("  3) 🔀 混合模式 - 文章 + 图片 + 视频")
        mode_choice = input("选择 [1/2/3] (默认: 1): ").strip() or '1'

        modes = {'1': 'image', '2': 'video', '3': 'hybrid'}
        mode = modes.get(mode_choice, 'image')

        # 2. 选择平台（支持多选）
        print("\n📱 请选择目标平台（可多选，如 1,2 或 12）:")
        print("  1) 📱 抖音/快手 - 1080×1920 (9:16)")
        print("  2) 🎬 B站/YouTube - 1920×1080 (16:9)")
        print("  3) 💬 视频号 - 1920×1080 (16:9)")
        print("  4) ⚙️  自定义")
        platform_input = input("选择 [1/2/3/4] (默认: 2): ").strip() or '2'

        # 解析多选（支持 "1,2" "1 2" "12"）
        selected = []
        for c in platform_input.replace(',', ' ').split():
            selected.extend(list(c.strip()))
        selected = list(dict.fromkeys(selected))  # 去重保序

        platforms = {
            '1': ('1080x1920', 'tiktok', 'pixelize', 30),
            '2': ('1920x1080', 'youtube', 'fade', 30),
            '3': ('1920x1080', 'youtube', 'fade', 30),
            '4': None
        }

        dual_version = False
        if '4' in selected:
            # 自定义
            resolution = input("分辨率 (如 1920x1080): ").strip() or '1920x1080'
            subtitle_style = input("字幕样式 [news/youtube/tiktok/minimal] (默认: news): ").strip() or 'news'
            transition = input("转场效果 [fade/wipeleft/wiperight/slideleft/pixelize/none] (默认: fade): ").strip() or 'fade'
            fps = int(input("帧率 (默认: 30): ").strip() or '30')
        else:
            # 取第一个有效平台为主配置
            first = next((s for s in selected if s in platforms and s != '4'), '2')
            resolution, subtitle_style, transition, fps = platforms[first]
            # 如果选了多个不同比例的平台，开启双版本
            unique_res = set()
            for s in selected:
                if s in platforms and platforms[s]:
                    unique_res.add(platforms[s][0])
            if len(unique_res) > 1:
                dual_version = True
                print(f"   🔀 已选多平台，将同时生成横竖双版本")

        # 3. 选择音色
        print("\n🎙️  请选择默认AI音色:")
        print("  1) 👩 晓晓 - 活泼女声 (推荐)")
        print("  2) 👨 云扬 - 成熟男声")
        print("  3) 👦 云希 - 年轻男声")
        print("  4) 👩 晓伊 - 成熟女声")
        print("  5) 👩 云夏 - 年轻女声")
        print("  6) 👨 云健 - 新闻播报男声")
        voice_choice = input("选择 [1-6] (默认: 1): ").strip() or '1'

        voices = {
            '1': 'Xiaoxiao', '2': 'Yunyang', '3': 'Yunxi',
            '4': 'Xiaoyi', '5': 'Yunxia', '6': 'Yunjian'
        }
        voice = voices.get(voice_choice, 'Xiaoxiao')

    # 4. 创建目录结构
    print(f"\n📂 创建目录结构...")
    dirs = ['01_article']
    if mode in ('image', 'hybrid'):
        dirs.append('03_images')
    if mode in ('video', 'hybrid'):
        dirs.append('04_videos')
    dirs.extend(['06_scenes', '07_final', '02_bgm', '02_sfx'])

    for d in dirs:
        (project_dir / d).mkdir(exist_ok=True)
        print(f"  ✅ {d}/")

    # 5. 创建示例文章
    article_path = project_dir / '01_article' / '文章.md'
    if not article_path.exists():
        article_templates = {
            'image': f"""# 示例文章

@全局:{'女声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '男声'}
@默认图: 01

@{'男声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '女声'}: 大家好，欢迎收看今天的节目。

这是第二段内容，使用默认图片。

@{'男声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '女声'}: 感谢观看，我们下期再见！
""",
            'video': f"""# 示例文章

@全局:{'女声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '男声'}

这是视频配音文案。

第一段内容。

第二段内容。

结尾部分。
""",
            'hybrid': f"""# 示例文章

@全局:{'女声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '男声'}
@默认图: 01

@{'男声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '女声'}: 大家好，欢迎收看。

这是第二段内容。

@{'男声' if voice in ['Xiaoxiao', 'Xiaoyi', 'Yunxia'] else '女声'}: 感谢观看！
"""
        }
        if template_config:
            article_path.write_text(template_config['article'], encoding='utf-8')
        else:
            article_path.write_text(article_templates[mode], encoding='utf-8')
        print(f"  ✅ 01_article/文章.md")

    # 6. 生成项目配置文件
    config_path = project_dir / '.video_config.json'
    config = {
        'mode': mode,
        'resolution': resolution,
        'fps': fps,
        'subtitle': True,
        'subtitle_style': subtitle_style,
        'subtitle_animation': 'none',
        'transition': transition,
        'voice': voice,
        'transition_duration': 0.5,
        'rate': '+18%',
        'scene_fade': 0.0,
        'bgm_volume': 0.25,
        'watermark': None,
        'watermark_position': 'bottom-right',
        'sfx': False,
        'dual_version': dual_version,
        'normalize_audio': False,
        'subtitle_mode': 'sentence',
        'subtitle_gap': 0.1,
        'intro_text': None,
        'outro_text': None,
        'created': str(datetime.datetime.now())
    }
    with open(config_path, 'w', encoding='utf-8') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)
    print(f"  ✅ .video_config.json")

    # 7. 生成参数说明文档
    param_doc_path = project_dir / '参数说明.md'
    if not param_doc_path.exists():
        param_doc = f"""# 项目参数说明

> 本文件由初始化向导自动生成，可直接编辑 `.video_config.json` 修改参数，
> 无需每次在命令行输入。

## 当前配置

```json
{json.dumps(config, ensure_ascii=False, indent=2)}
```

## 参数详解

### 基础参数

| 参数 | 当前值 | 中文说明 |
|------|--------|----------|
| `mode` | `{mode}` | 项目模式: `image`(图片+配音) / `video`(视频+配音) / `hybrid`(混合) |
| `resolution` | `{resolution}` | 输出分辨率，如 `1920x1080`(横屏) / `1080x1920`(竖屏) |
| `fps` | `{fps}` | 帧率，默认 30 |
| `voice` | `{voice}` | 默认AI配音音色 |
| `rate` | `+18%` | 语速调节，如 `+10%`(加快) / `-10%`(放慢) |

### 视觉参数

| 参数 | 当前值 | 中文说明 |
|------|--------|----------|
| `subtitle` | `true` | 是否启用字幕，`true`(开启) / `false`(关闭) |
| `subtitle_style` | `{subtitle_style}` | 字幕样式预设，见下表 |
| `subtitle_animation` | `none` | 字幕动画: `none`(无) / `slide_up`(上滑进入) / `fade_in`(淡入) |
| `transition` | `{transition}` | 场景转场效果，见下表 |
| `transition_duration` | `0.5` | 转场时长(秒)，建议 0.3~1.0 |
| `scene_fade` | `0.0` | 场景淡入淡出(秒)，建议 0.2~0.5，0 为关闭 |
| `watermark` | `null` | 水印图片路径，如 `logo.png`，`null` 为不添加 |
| `watermark_position` | `bottom-right` | 水印位置: `top-left`(左上) / `top-right`(右上) / `bottom-left`(左下) / `bottom-right`(右下) / `center`(居中) |

### 音频参数

| 参数 | 当前值 | 中文说明 |
|------|--------|----------|
| `bgm_volume` | `0.25` | 背景音乐音量，范围 0.0~1.0，0 为静音 |
| `sfx` | `false` | 是否启用转场音效，`true`(开启，自动读取 `02_sfx/` 目录) / `false`(关闭) |
| `dual_version` | `{dual_version}` | 是否同时生成横竖双版本，`true`(开启) / `false`(关闭) |

## 字幕样式列表

| 样式名 | 中文名 | 特点 |
|--------|--------|------|
| `news` | 新闻 | 底部黑底白字，带边框，适合新闻/解说 |
| `youtube` | YouTube | 底部黄色大字，黑色描边，适合英文/教程 |
| `tiktok` | 抖音 | 居中白色大字，半透黑底，适合竖屏短视频 |
| `minimal` | 极简 | 细边框，无背景框，适合文艺/纪录片风格 |

## 转场效果列表

| 效果名 | 中文名 | 效果名 | 中文名 |
|--------|--------|--------|--------|
| `fade` | 淡入淡出 | `dissolve` | 溶解 |
| `wipeleft` | 向左擦除 | `wiperight` | 向右擦除 |
| `wipeup` | 向上擦除 | `wipedown` | 向下擦除 |
| `slideleft` | 向左滑动 | `slideright` | 向右滑动 |
| `slideup` | 向上滑动 | `slidedown` | 向下滑动 |
| `smoothleft` | 平滑向左 | `smoothright` | 平滑向右 |
| `smoothup` | 平滑向上 | `smoothdown` | 平滑向下 |
| `circlecrop` | 圆形裁剪 | `rectcrop` | 矩形裁剪 |
| `circleclose` | 圆形收缩 | `circleopen` | 圆形展开 |
| `horzclose` | 水平收缩 | `horzopen` | 水平展开 |
| `vertclose` | 垂直收缩 | `vertopen` | 垂直展开 |
| `pixelize` | 像素化 | `radial` | 径向扩散 |
| `distance` | 距离模糊 | `fadeblack` | 黑场淡出 |
| `fadewhite` | 白场淡出 | `diagtl` | 对角线(左上) |
| `diagtr` | 对角线(右上) | `diagbl` | 对角线(左下) |
| `diagbr` | 对角线(右下) | `hlslice` | 水平左切片 |
| `hrslice` | 水平右切片 | `vuslice` | 垂直上切片 |
| `vdslice` | 垂直下切片 | `none` | 无转场(直接切换) |

## 音色列表

| 音色名 | 中文名 | 特点 |
|--------|--------|------|
| `Xiaoxiao` | 晓晓 | 活泼温暖女声（默认女声） |
| `Xiaoyi` | 晓伊 | 成熟稳重女声 |
| `Yunxia` | 云夏 | 年轻女声 |
| `Yunyang` | 云扬 | 成熟稳重男声（默认男声） |
| `Yunxi` | 云希 | 年轻男声 |
| `Yunjian` | 云健 | 新闻播报男声 |

## 快速修改方法

### 方法1：直接改配置文件（推荐）
编辑 `.video_config.json`，修改对应值，下次运行即可生效：

```bash
python3 video_generator_pro.py -p {project_dir}
```

### 方法2：命令行覆盖
命令行参数优先级高于配置文件，适合临时调试：

```bash
python3 video_generator_pro.py -p {project_dir} \
    --scene-fade 0.3 \
    --subtitle --subtitle-animation slide_up \
    --transition pixelize \
    --watermark logo.png \
    --watermark-position bottom-right \
    --sfx
```

---
*生成时间: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*
"""
        param_doc_path.write_text(param_doc, encoding='utf-8')
        print(f"  ✅ 参数说明.md")

    # 7. 显示后续步骤
    print(f"\n{'='*60}")
    print("✅ 项目初始化完成!")
    print(f"{'='*60}")
    print(f"📁 项目路径: {project_dir}")
    print(f"\n📋 后续步骤:")
    if mode in ('image', 'hybrid'):
        print(f"  1. 放入图片到: {project_dir}/03_images/")
        print(f"     命名: 01.jpg, 02.jpg, ... 或 吃饭.jpg, 睡觉.png, ...")
    if mode in ('video', 'hybrid'):
        print(f"  1. 放入视频到: {project_dir}/04_videos/")
        print(f"     命名: scene_01.mp4, scene_02.mp4, ...")
    print(f"  2. 编辑文章: {project_dir}/01_article/文章.md")
    print(f"  3. (可选) 放入 BGM 到: {project_dir}/02_bgm/")
    print(f"     支持: .mp3 .wav .aac .m4a，自动循环播放")
    print(f"  4. (可选) 放入转场音效到: {project_dir}/02_sfx/")
    print(f"     支持: .mp3 .wav，自动在转场时混入")
    print(f"  5. 生成视频:")
    print(f"     python3 video_generator_pro.py -p {project_dir} --voice {voice}")

    if mode in ('image', 'hybrid'):
        print(f"\n💡 提示: 需要 {resolution} 比例的图片素材")

    return True


def merge_project_config(args, project_dir: Path):
    """加载项目 .video_config.json，命令行参数优先级 > 配置文件 > argparse 默认值"""
    config_path = project_dir / '.video_config.json'
    if not config_path.exists():
        return args

    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
    except Exception:
        return args

    # 解析命令行中显式指定的参数（长参数和短参数）
    explicit = set()
    # 短参数到长参数的映射
    short_to_long = {
        's': 'subtitle', 'p': 'project', 't': 'transition',
        'b': 'bgm', 'o': 'output'
    }
    i = 0
    argv = sys.argv[1:]
    while i < len(argv):
        arg = argv[i]
        if arg.startswith('--'):
            # --key=value 或 --key
            key = arg[2:].split('=')[0].replace('-', '_')
            explicit.add(key)
            # --no-xxx 也视为 xxx 被显式指定
            if key.startswith('no_') and key[3:] in config_keys:
                explicit.add(key[3:])
        elif arg.startswith('-') and len(arg) == 2:
            # 短参数 -k，映射为长参数名
            short_key = arg[1:]
            explicit.add(short_to_long.get(short_key, short_key))
            if i + 1 < len(argv) and not argv[i + 1].startswith('-'):
                i += 1
        i += 1

    # 支持的配置键映射（配置文件键名 -> args 属性名）
    config_keys = [
        'voice', 'resolution', 'fps', 'subtitle_style', 'transition',
        'transition_duration', 'rate', 'bgm_volume', 'subtitle', 'output',
        'scene_fade', 'watermark', 'watermark_position', 'sfx', 'subtitle_animation',
        'dual_version', 'normalize_audio', 'subtitle_mode', 'intro_text', 'outro_text'
    ]

    applied = []
    for key in config_keys:
        if key in config and key not in explicit:
            if hasattr(args, key):
                old_val = getattr(args, key)
                new_val = config[key]
                if old_val != new_val:
                    setattr(args, key, new_val)
                    applied.append(f"{key}={new_val}")

    if applied:
        print(f"   📋 已加载项目配置: {', '.join(applied)}")

    return args


def check_project_materials(project_dir: Path, image_assignments: list = None) -> dict:
    """检查项目素材完整性

    Returns:
        dict: 检查结果
    """
    print(f"\n{'='*60}")
    print("🔍 素材检查报告")
    print(f"{'='*60}")

    result = {
        'project_dir': str(project_dir),
        'valid': True,
        'warnings': [],
        'errors': [],
        'stats': {}
    }

    # 1. 检查文章
    article_dir = project_dir / '01_article'
    article_files = list(article_dir.glob('*.md')) + list(article_dir.glob('*.txt')) if article_dir.exists() else []
    if article_files:
        result['stats']['articles'] = len(article_files)
        print(f"  ✅ 文章: {len(article_files)} 个")
        for f in article_files:
            print(f"     - {f.name}")
    else:
        result['warnings'].append("未找到文章文件（01_article/*.md 或 *.txt）")
        print(f"  ⚠️  未找到文章文件")

    # 2. 检查音频
    audio_dir = project_dir / '05_audio'
    audio_files = sorted([f for f in audio_dir.iterdir() if f.suffix == '.mp3']) if audio_dir.exists() else []
    result['stats']['audio_files'] = len(audio_files)
    if audio_files:
        total_duration = sum(get_media_duration(str(f)) for f in audio_files)
        print(f"  ✅ 音频: {len(audio_files)} 个 (总时长: {total_duration:.1f}s)")
    else:
        if article_files:
            result['warnings'].append("音频将自动生成（运行时）")
            print(f"  ⏳ 音频: 未找到，将自动生成")
        else:
            result['errors'].append("无音频且无文章，无法生成")
            print(f"  ❌ 音频: 无音频且无文章")
            result['valid'] = False

    # 3. 检查图片/视频素材
    images_dir = project_dir / '03_images'
    videos_dir = project_dir / '04_videos'

    images = []
    videos = []

    if images_dir.exists():
        images = sorted([f for f in images_dir.iterdir()
                        if f.suffix.lower() in ['.jpg', '.jpeg', '.png', '.webp', '.gif', '.heic']])
    if videos_dir.exists():
        videos = sorted([f for f in videos_dir.iterdir()
                        if f.suffix.lower() in ['.mp4', '.mov', '.avi', '.mkv']])

    result['stats']['images'] = len(images)
    result['stats']['videos'] = len(videos)

    if videos:
        print(f"  ✅ 视频: {len(videos)} 个")
        for v in videos:
            duration = get_media_duration(str(v))
            print(f"     - {v.name} ({duration:.1f}s)")
    if images:
        print(f"  ✅ 图片: {len(images)} 张")
        for img in images:
            print(f"     - {img.name}")

    # 4. 检查场景片段
    scenes_dir = project_dir / '06_scenes'
    existing_scenes = sorted([f for f in scenes_dir.iterdir() if f.suffix == '.mp4']) if scenes_dir.exists() else []
    result['stats']['existing_scenes'] = len(existing_scenes)
    if existing_scenes:
        print(f"  ✅ 场景片段: {len(existing_scenes)} 个 (将跳过生成)")

    # 5. 检查 BGM
    bgm_dir = project_dir / '02_bgm'
    bgm_files = []
    if bgm_dir.exists():
        music_exts = ['.mp3', '.wav', '.aac', '.flac', '.m4a', '.ogg']
        bgm_files = sorted([f for f in bgm_dir.iterdir() if f.suffix.lower() in music_exts])
    result['stats']['bgm'] = len(bgm_files)
    if bgm_files:
        print(f"  ✅ BGM: {len(bgm_files)} 个")
        for bgm in bgm_files:
            duration = get_media_duration(str(bgm))
            print(f"     - {bgm.name} ({duration:.1f}s)")

    # 6. 匹配检查
    if audio_files:
        audio_count = len(audio_files)
        if videos:
            # 视频模式
            video_count = len(videos)
            if audio_count != video_count:
                result['warnings'].append(f"音频({audio_count})和视频({video_count})数量不匹配")
                print(f"  ⚠️  音频({audio_count})和视频({video_count})数量不匹配")
        elif images:
            # 图片模式
            img_count = len(images)
            if audio_count > img_count:
                result['warnings'].append(f"音频({audio_count})多于图片({img_count})，部分场景将复用图片")
                print(f"  ⚠️  音频({audio_count})多于图片({img_count})，部分场景将复用图片")

    # 6. 加载项目配置
    config_path = project_dir / '.video_config.json'
    if config_path.exists():
        with open(config_path, 'r') as f:
            config = json.load(f)
        print(f"\n📋 项目配置:")
        print(f"  模式: {config.get('mode', '未设置')}")
        print(f"  分辨率: {config.get('resolution', '1920x1080')}")
        print(f"  帧率: {config.get('fps', 30)}")
        print(f"  字幕: {config.get('subtitle_style', 'news')}")
        print(f"  转场: {config.get('transition', 'fade')}")
        print(f"  音色: {config.get('voice', 'Xiaoxiao')}")

    print(f"\n{'─'*60}")
    if result['errors']:
        print(f"❌ 发现 {len(result['errors'])} 个错误，无法生成")
        for e in result['errors']:
            print(f"   - {e}")
        result['valid'] = False
    elif result['warnings']:
        print(f"⚠️  发现 {len(result['warnings'])} 个警告，但可以生成")
        for w in result['warnings']:
            print(f"   - {w}")
    else:
        print(f"✅ 素材检查通过，可以生成")
    print(f"{'─'*60}")

    return result


def main():
    parser = argparse.ArgumentParser(
        description="视频合成工具 Pro 版",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
项目文件夹结构:
  01_article/  - 文章文件 (.md 或 .txt)，自动生成音频
  03_images/   - 图片素材 (.jpg/.png)，配合音频生成视频
  04_videos/   - 视频素材 (.mp4/.mov)，配合音频替换原声
  05_audio/    - 音频文件 (.mp3)，手动放置或自动生成
  04_final/    - 最终输出视频

字幕样式 (--subtitle-style):
  news     - 新闻风格 (白字黑底)
  youtube  - YouTube风格 (黄字黑边)
  minimal  - 极简风格 (白字)
  tiktok   - 抖音风格 (大字居中)

转场效果 (--transition):
  fade        - 淡入淡出
  wipeleft    - 向左擦除
  wiperight   - 向右擦除
  slideleft   - 向左滑动
  slideright  - 向右滑动
  pixelize    - 像素化
  none        - 无转场

AI配音音色 (--voice):
  Xiaoxiao  - 晓晓女声 (默认)
  Yunyang   - 云扬男声
  Yunxi     - 云希年轻男声
  Xiaoyi    - 晓伊成熟女声

示例:
  # 基础用法 (自动识别图片/视频+音频)
  python3 video_generator_pro.py -p projects/XXX

  # 从文章自动生成音频并合成
  python3 video_generator_pro.py -p projects/XXX --voice Xiaoxiao --rate +0%

  # 快速预览（只生成第一个场景，低画质，快速验证）
  python3 video_generator_pro.py -p projects/XXX --preview

  # 完整功能
  python3 video_generator_pro.py -p projects/XXX \\
    --subtitle-style news \\
    --transition fade \\
    --intro intro.mp4 --outro outro.mp4 \\
    --bgm music.mp3 --bgm-volume 0.2

  # 批量处理（自动发现当前目录项目）
  python3 video_generator_pro.py --batch

  # 队列处理（指定多个项目，失败自动继续）
  python3 video_generator_pro.py --queue projects/A projects/B projects/C
        """
    )

    parser.add_argument('--project', '-p', help='项目路径')
    parser.add_argument('--batch', action='store_true', help='批量处理多个项目')
    parser.add_argument('--subtitle', '-s', action='store_true', default=True,
                       help='添加字幕（默认开启）')
    parser.add_argument('--no-subtitle', action='store_false', dest='subtitle',
                       help='关闭字幕')
    parser.add_argument('--subtitle-style', default=None,
                       choices=list(SUBTITLE_STYLES.keys()),
                       help='字幕样式 (未指定时自动适配: 横屏 news, 竖屏 tiktok)')
    parser.add_argument('--transition', '-t', default='fade',
                       choices=list(TRANSITIONS.keys()),
                       help='转场效果 (默认: fade)')
    parser.add_argument('--transition-duration', type=float, default=0.5,
                       help='转场时长 (秒, 默认: 0.5)')
    parser.add_argument('--intro', help='片头视频路径')
    parser.add_argument('--outro', help='片尾视频路径')
    parser.add_argument('--bgm', '-b', help='背景音乐路径')
    parser.add_argument('--bgm-volume', type=float, default=0.25,
                       help='背景音乐音量 (0.0-1.0, 默认: 0.25)')
    parser.add_argument('--resolution', default='1920x1080',
                       help='分辨率 (默认: 1920x1080)')
    parser.add_argument('--fps', type=int, default=30, help='帧率 (默认: 30)')
    parser.add_argument('--output', '-o', help='输出文件名')
    parser.add_argument('--voice', default='Xiaoxiao',
                       choices=['Xiaoxiao', 'Xiaoyi', 'Yunxi', 'Yunjian', 'Yunxia', 'Yunyang',
                               'HsiaoChen', 'HsiaoYu', 'YunJhe', 'HiuMaan', 'HiuGaai', 'WanLung'],
                       help='AI配音音色 (默认: Xiaoxiao)')
    parser.add_argument('--rate', default='+18%',
                       help='语速调节 (默认: +18%%)')
    parser.add_argument('--init', action='store_true',
                       help='交互式初始化新项目')
    parser.add_argument('--template',
                       choices=list(PROJECT_TEMPLATES.keys()),
                       help='使用预设模板初始化 (news/food/tutorial/education)')
    parser.add_argument('--check', action='store_true',
                       help='只检查素材，不生成视频')
    parser.add_argument('--no-parallel', action='store_true',
                       help='关闭并行生成（默认开启，最多3并发）')
    parser.add_argument('--preview', action='store_true',
                       help='预览模式：只生成第一个场景，低分辨率快速验证效果')
    parser.add_argument('--preview-duration', type=float, default=5.0,
                       help='预览模式限制单场景时长（默认5秒）')
    parser.add_argument('--skip-pre-check', action='store_true',
                       help='跳过素材预检（默认执行预检）')
    parser.add_argument('--regenerate-audio', action='store_true',
                       help='强制重新生成音频（保留已有场景）')
    parser.add_argument('--normalize-audio', action='store_true',
                       help='音频响度标准化到 -14 LUFS（YouTube标准）')
    parser.add_argument('--intro-text', help='片头文字（自动生成片头视频）')
    parser.add_argument('--outro-text', help='片尾文字（自动生成片尾视频）')
    parser.add_argument('--generate-copy', action='store_true',
                       help='根据文章内容自动生成多平台发布文案（需配置 LLM API）')
    parser.add_argument('--llm-provider', default='kimi',
                       choices=['kimi', 'deepseek'],
                       help='LLM 提供商 (默认: kimi)')
    parser.add_argument('--llm-api-key', help='LLM API 密钥（也可通过 OPENAI_API_KEY/KIMI_API_KEY/DEEPSEEK_API_KEY 环境变量设置）')
    parser.add_argument('--llm-base-url', help='LLM API 基础地址（默认由 --llm-provider 决定）')
    parser.add_argument('--llm-model', help='LLM 模型名称（默认由 --llm-provider 决定）')
    parser.add_argument('--import-ppt', metavar='PATH',
                       help='导入PPT/Keynote：提取图片和备注文本生成新项目')
    parser.add_argument('--queue', nargs='+', metavar='PATH',
                       help='队列模式：顺序处理多个项目，失败自动继续下一个')
    parser.add_argument('--scene-fade', type=float, default=0.0,
                       help='场景淡入淡出时长 (秒, 默认: 0.0, 建议 0.3-0.5)')
    parser.add_argument('--watermark', help='水印图片路径')
    parser.add_argument('--watermark-position', default='bottom-right',
                       choices=['top-left', 'top-right', 'bottom-left', 'bottom-right', 'center'],
                       help='水印位置 (默认: bottom-right)')
    parser.add_argument('--sfx', action='store_true',
                       help='启用转场音效 (自动查找项目 02_sfx/ 目录)')
    parser.add_argument('--subtitle-animation', default='none',
                       choices=['none', 'slide_up', 'fade_in'],
                       help='字幕动画效果 (默认: none)')
    parser.add_argument('--subtitle-mode', default='sentence',
                       choices=['full', 'sentence'],
                       help='字幕显示模式: full(整段显示) / sentence(逐句显示) (默认: sentence)')
    parser.add_argument('--subtitle-gap', type=float, default=0.1,
                       help='逐句字幕句间黑屏间隔(秒)。0=无缝衔接下一句,>0=读完后黑屏N秒再出下一句 (默认: 0.1)')
    parser.add_argument('--dual-version', action='store_true',
                       help='同时生成横竖双版本（如当前为横屏则额外生成竖屏，反之亦然）')
    parser.add_argument('--batch-variants-dir', metavar='PATH',
                       help='批量模板生成：扫描目录下的子目录作为素材变体，每套素材+模板文章生成一个视频（矩阵号）')
    parser.add_argument('--whisper-transcribe', action='store_true',
                       help='语音识别：使用本地 faster-whisper 识别视频语音，自动生成文章和字幕（视频模式专用）')
    parser.add_argument('--auto-article', metavar='TITLE',
                       help='AI 自动生成文章：输入标题，调用 LLM 直接生成口播文章（需要配置 API Key）')
    parser.add_argument('--search-web', action='store_true',
                       help='启用联网搜索：让 LLM 获取实时数据生成文章（需要 API 支持 web_search）')
    parser.add_argument('--auto-images', action='store_true',
                       help='AI 自动配图：根据文章内容自动搜索/生成图片并插入到文章中')
    parser.add_argument('--image-provider', default='pollinations',
                       choices=['pollinations', 'unsplash', 'pexels'],
                       help='图片提供商 (默认: pollinations，免费免 key)')
    parser.add_argument('--image-api-key', help='图片 API 密钥（Unsplash/Pexels 需要，Pollinations 不需要）')

    args = parser.parse_args()

    # 检查 ffmpeg
    if not shutil.which('ffmpeg'):
        print("❌ 未找到 ffmpeg，请先安装: brew install ffmpeg")
        sys.exit(1)

    # 项目初始化向导
    if args.init or args.template:
        if args.project:
            project_dir = Path(args.project)
        else:
            project_dir = Path(input("请输入项目路径: ").strip().replace('\r', ''))
        if init_project_wizard(project_dir, template=args.template):
            if input("\n是否立即检查素材? [y/N]: ").strip().replace('\r', '').lower() == 'y':
                check_project_materials(project_dir)
        sys.exit(0)

    # AI 自动生成文章
    auto_images_done = False
    if args.auto_article:
        title = args.auto_article.strip()
        if args.project:
            project_dir = Path(args.project)
        else:
            import re
            safe_name = re.sub(r'[^\w\u4e00-\u9fff]+', '_', title)[:30].strip('_')
            if not safe_name:
                safe_name = 'auto_article'
            project_dir = Path('projects') / safe_name

        project_dir.mkdir(parents=True, exist_ok=True)
        # 确保标准目录结构存在
        for subdir in ['01_article', '02_bgm', '03_images']:
            (project_dir / subdir).mkdir(exist_ok=True)

        # 智能判断是否需要联网搜索
        search_web = getattr(args, 'search_web', False)
        if not search_web and _needs_realtime_search(title):
            search_web = True
            print(f"   🧠 检测到时效性关键词，自动开启联网搜索")

        article_path = auto_generate_article_from_title(
            title,
            project_dir,
            api_key=getattr(args, 'llm_api_key', None),
            base_url=getattr(args, 'llm_base_url', None),
            model=getattr(args, 'llm_model', None),
            provider=getattr(args, 'llm_provider', 'kimi'),
            search_web=search_web
        )
        if not article_path:
            sys.exit(1)

        # 如果同时指定了 --auto-images，执行配图并继续生成视频
        if getattr(args, 'auto_images', False):
            args.project = str(project_dir)
            print(f"\n{'='*60}")
            print("🎨 第二阶段：自动配图...")
            print(f"{'='*60}")
            img_count = auto_generate_images_for_project(
                project_dir,
                image_provider=getattr(args, 'image_provider', 'pollinations'),
                image_api_key=getattr(args, 'image_api_key', None),
                llm_provider=getattr(args, 'llm_provider', 'kimi'),
                llm_api_key=getattr(args, 'llm_api_key', None),
                llm_base_url=getattr(args, 'llm_base_url', None),
                llm_model=getattr(args, 'llm_model', None)
            )
            auto_images_done = True
            if img_count == 0:
                print("⚠️  自动配图未成功，继续生成视频（可能使用默认图片）")
            # 继续执行，不退出（会进入 elif args.project 生成视频）
        else:
            print(f"\n{'='*60}")
            print("🎉 第一阶段完成！")
            print(f"{'='*60}")
            print(f"   项目路径: {project_dir}")
            print(f"   文章文件: {project_dir}/01_article/文章.md")
            print(f"\n   📋 下一步:")
            print(f"   1. 放入图片到 {project_dir}/03_images/")
            print(f"   2. 放入 BGM 到 {project_dir}/02_bgm/")
            print(f"   3. 生成视频:")
            print(f"      python3 01_核心脚本/video_generator_pro.py -p {project_dir}")
            print(f"{'='*60}")
            sys.exit(0)

    # PPT导入
    if args.import_ppt:
        ppt_path = Path(args.import_ppt)
        if args.project:
            project_dir = Path(args.project)
        else:
            project_dir = Path('projects') / ppt_path.stem
        import_ppt_project(ppt_path, project_dir)
        sys.exit(0)

    # 批量模板生成（矩阵号）
    if args.batch_variants_dir:
        variants_dir = Path(args.batch_variants_dir)
        if not variants_dir.exists() or not variants_dir.is_dir():
            print(f"❌ 变体目录不存在: {variants_dir}")
            sys.exit(1)

        template_dir = Path(args.project) if args.project else Path('.')
        if not template_dir.exists():
            print(f"❌ 模板项目不存在，请用 -p 指定")
            sys.exit(1)

        # 扫描变体目录
        variant_dirs = sorted([d for d in variants_dir.iterdir() if d.is_dir()])
        if not variant_dirs:
            print(f"❌ 变体目录下没有找到子目录: {variants_dir}")
            sys.exit(1)

        print(f"\n{'='*60}")
        print(f"🔄 批量模板生成: {template_dir.name}")
        print(f"   模板项目: {template_dir}")
        print(f"   变体目录: {variants_dir}")
        print(f"   发现 {len(variant_dirs)} 套素材")
        print(f"{'='*60}")

        batch_start = time.time()
        results = []
        success_count = 0
        fail_count = 0

        for i, variant in enumerate(variant_dirs, 1):
            print(f"\n📌 [{i}/{len(variant_dirs)}] 变体: {variant.name}")
            print(f"{'─'*60}")

            # 创建临时项目
            temp_dir = Path(tempfile.mkdtemp(prefix=f"batch_{template_dir.name}_"))
            try:
                # 复制模板结构
                for subdir in ['01_article', '02_bgm', '02_sfx', '03_images', '04_videos']:
                    src = template_dir / subdir
                    dst = temp_dir / subdir
                    if src.exists():
                        if subdir == '03_images' or subdir == '02_bgm':
                            # 图片和BGM从变体复制
                            variant_src = variant / subdir
                            if variant_src.exists():
                                shutil.copytree(variant_src, dst)
                            elif src.exists():
                                shutil.copytree(src, dst)
                        else:
                            shutil.copytree(src, dst)

                # 复制配置文件
                config_src = template_dir / '.video_config.json'
                if config_src.exists():
                    shutil.copy(config_src, temp_dir / '.video_config.json')

                # 复制水印
                wm_src = template_dir / 'watermark.png'
                if wm_src.exists():
                    shutil.copy(wm_src, temp_dir / 'watermark.png')

                # 清理旧音频和场景（强制重新生成，因为图片变了）
                for d in ['05_audio', '06_scenes']:
                    old = temp_dir / d
                    if old.exists():
                        shutil.rmtree(old)

                # 加载配置并生成
                merge_project_config(args, temp_dir)
                result = process_project(temp_dir, args)

                if result:
                    # 复制最终视频到模板项目的 07_final/ 下，带变体名
                    final_dir = template_dir / '07_final'
                    final_dir.mkdir(exist_ok=True)
                    dest_name = f"{template_dir.name}_{variant.name}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.mp4"
                    dest_path = final_dir / dest_name
                    shutil.copy(str(result), str(dest_path))
                    print(f"   ✅ 已保存: {dest_path.name}")
                    results.append((variant.name, True, str(dest_path)))
                    success_count += 1
                else:
                    results.append((variant.name, False, None))
                    fail_count += 1

            except Exception as e:
                print(f"\n❌ 变体异常: {e}")
                traceback.print_exc()
                results.append((variant.name, False, None))
                fail_count += 1
            finally:
                # 清理临时目录
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)

        batch_total = time.time() - batch_start

        # 汇总报告
        print(f"\n{'='*60}")
        print(f"📊 批量模板生成完成报告")
        print(f"{'='*60}")
        print(f"  📁 总变体: {len(variant_dirs)} 个")
        print(f"  ✅ 成功:   {success_count} 个")
        print(f"  ❌ 失败:   {fail_count} 个")
        print(f"  ⏱️  总耗时:  {batch_total:.1f}s")
        if success_count > 0:
            print(f"  ⏱️  平均耗时: {batch_total / success_count:.1f}s/变体")
        print(f"\n  📋 明细:")
        for name, ok, path in results:
            status = "✅" if ok else "❌"
            detail = f" → {Path(path).name}" if ok else ""
            print(f"    {status} {name}{detail}")
        print(f"{'='*60}")
        sys.exit(0)

    # 批量处理 / 队列处理
    if args.batch or args.queue:
        if args.queue:
            # 队列模式：从命令行获取指定项目列表
            queue_paths = [Path(p) for p in args.queue]
            projects = []
            for p in queue_paths:
                if p.exists() and p.is_dir():
                    projects.append(p)
                else:
                    print(f"⚠️  跳过不存在的项目: {p}")
        else:
            # 批量模式：自动发现当前目录下的项目
            projects = [p for p in Path('.').iterdir() if p.is_dir() and ((p / '04_videos').exists() or (p / '05_audio').exists())]

        if not projects:
            print("❌ 未找到项目文件夹")
            sys.exit(1)

        mode_label = "队列" if args.queue else "批量"
        print(f"\n{'='*60}")
        print(f"🔄 {mode_label}处理 {len(projects)} 个项目...")
        print(f"{'='*60}")

        queue_start = time.time()
        results = []
        success_count = 0
        fail_count = 0
        skip_count = 0

        for i, project in enumerate(projects, 1):
            print(f"\n📌 [{i}/{len(projects)}] 项目: {project.name}")
            print(f"{'─'*60}")

            # 加载项目配置
            merge_project_config(args, project)

            try:
                result = process_project(project, args)
                if result:
                    results.append((project.name, True, str(result)))
                    success_count += 1
                else:
                    results.append((project.name, False, None))
                    fail_count += 1
            except KeyboardInterrupt:
                print(f"\n⛔ 用户中断，终止{mode_label}处理")
                sys.exit(1)
            except Exception as e:
                print(f"\n❌ 项目异常: {e}")
                traceback.print_exc()
                results.append((project.name, False, None))
                fail_count += 1

        queue_total = time.time() - queue_start

        # 汇总报告
        print(f"\n{'='*60}")
        print(f"📊 {mode_label}处理完成报告")
        print(f"{'='*60}")
        print(f"  📁 总项目: {len(projects)} 个")
        print(f"  ✅ 成功:   {success_count} 个")
        print(f"  ❌ 失败:   {fail_count} 个")
        print(f"  ⏱️  总耗时:  {queue_total:.1f}s")
        if success_count > 0:
            print(f"  ⏱️  平均耗时: {queue_total / success_count:.1f}s/项目")
        print(f"\n  📋 项目明细:")
        for name, ok, path in results:
            status = "✅" if ok else "❌"
            detail = f" → {path}" if ok else ""
            print(f"    {status} {name}{detail}")
        print(f"{'='*60}")

    elif args.project:
        # 单个项目
        project_dir = Path(args.project)
        if not project_dir.exists():
            print(f"❌ 项目不存在: {args.project}")
            sys.exit(1)

        # 加载项目配置（命令行参数优先级 > 配置文件）
        merge_project_config(args, project_dir)

        # 自动配图（如果指定且未在 --auto-article 中执行过）
        if getattr(args, 'auto_images', False) and not auto_images_done:
            img_count = auto_generate_images_for_project(
                project_dir,
                image_provider=getattr(args, 'image_provider', 'pollinations'),
                image_api_key=getattr(args, 'image_api_key', None),
                llm_provider=getattr(args, 'llm_provider', 'kimi'),
                llm_api_key=getattr(args, 'llm_api_key', None),
                llm_base_url=getattr(args, 'llm_base_url', None),
                llm_model=getattr(args, 'llm_model', None)
            )
            if img_count == 0:
                print("⚠️  自动配图未成功，继续生成视频（可能使用默认图片）")

        # 素材检查
        if args.check:
            check_project_materials(project_dir)
            sys.exit(0)

        # 生成发布文案（如果指定）
        if args.generate_copy:
            generate_publish_copy(
                project_dir,
                api_key=getattr(args, 'llm_api_key', None),
                base_url=getattr(args, 'llm_base_url', None),
                model=getattr(args, 'llm_model', None),
                provider=getattr(args, 'llm_provider', 'kimi')
            )

        result = process_project(project_dir, args)

        if result:
            print(f"\n▶️  播放命令: open '{result}'")
        else:
            sys.exit(1)

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == '__main__':
    main()
